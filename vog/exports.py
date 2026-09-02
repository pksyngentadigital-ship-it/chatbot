"""CSV / XLSX / PPTX generation, with no pandas.

pandas was pulled in for exactly two things — `DataFrame.to_csv` and
`DataFrame.to_excel` — and brought numpy with it. Together that is ~95MB,
which is most of a serverless bundle. The csv module and openpyxl (which
was already a dependency, since pandas used it as the Excel engine) do
the same job directly.
"""

import csv
import io

from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def _columns_of(rows: list[dict], columns: list[str] | None) -> list[str]:
    if columns:
        return columns
    seen: list[str] = []
    for row in rows:
        for key in row:
            if key not in seen:
                seen.append(key)
    return seen


def build_csv_bytes(rows: list[dict], columns: list[str] | None = None) -> bytes:
    if not rows:
        return b""
    cols = _columns_of(rows, columns)
    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=cols, extrasaction="ignore", lineterminator="\n")
    writer.writeheader()
    for row in rows:
        writer.writerow({c: row.get(c, "") for c in cols})
    return buf.getvalue().encode("utf-8")


def build_excel_bytes(rows: list[dict], columns: list[str] | None = None) -> bytes:
    if not rows:
        return b""
    cols = _columns_of(rows, columns)

    wb = Workbook()
    ws = wb.active
    ws.title = "Export"
    ws.append(cols)
    for cell in ws[1]:
        cell.font = Font(bold=True)

    for row in rows:
        ws.append([_excel_safe(row.get(c, "")) for c in cols])

    # Rough autofit — openpyxl has no real one, and unreadable 8-character
    # columns are the first thing anyone notices when opening the file.
    for i, col in enumerate(cols, start=1):
        widest = max([len(str(col))] + [len(str(r.get(col, ""))) for r in rows[:200]])
        ws.column_dimensions[get_column_letter(i)].width = min(max(widest + 2, 10), 60)
    ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _excel_safe(value):
    """Stop Excel interpreting a leading =, +, - or @ as a formula."""
    if isinstance(value, str) and value[:1] in ("=", "+", "-", "@"):
        return "'" + value
    return value


def build_pptx_bytes(
    title: str,
    subtitle: str,
    summary_lines: list[str],
    kpis: dict,
    chart_title: str = "",
    chart_labels: list | None = None,
    chart_values: list | None = None,
    chart_type: str = "column",
    table_headers: list | None = None,
    table_rows: list | None = None,
) -> bytes:
    """Title -> summary -> KPIs -> chart -> table. Native PowerPoint chart
    and table shapes, so the deck stays editable rather than being a
    picture of one."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    title_layout, bullet_layout, blank_layout = (
        prs.slide_layouts[0], prs.slide_layouts[1], prs.slide_layouts[6]
    )

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    _bullet_slide(prs, bullet_layout, "Summary", summary_lines or ["No summary available."], size=18)

    if kpis:
        _bullet_slide(prs, bullet_layout, "Key figures",
                      [f"{k}: {v}" for k, v in kpis.items()], size=20, bold=True)

    if chart_labels and chart_values:
        slide = prs.slides.add_slide(blank_layout)
        _heading(slide, chart_title or "Chart")
        data = CategoryChartData()
        data.categories = list(chart_labels)
        data.add_series(chart_title or "Value", list(chart_values))
        kind = XL_CHART_TYPE.LINE_MARKERS if chart_type == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
        slide.shapes.add_chart(kind, Inches(0.75), Inches(1.0), Inches(11.8), Inches(6.0), data)

    if table_headers and table_rows:
        slide = prs.slides.add_slide(blank_layout)
        _heading(slide, "Supporting data")
        shown = list(table_rows)[:15]
        table = slide.shapes.add_table(
            len(shown) + 1, len(table_headers),
            Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0)
        ).table
        for c, header in enumerate(table_headers):
            table.cell(0, c).text = str(header)
        for r, row in enumerate(shown, start=1):
            for c, value in enumerate(row):
                table.cell(r, c).text = str(value)[:300]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _heading(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    box.text_frame.text = text
    box.text_frame.paragraphs[0].font.size = Pt(24)
    box.text_frame.paragraphs[0].font.bold = True


def _bullet_slide(prs, layout, title: str, lines: list[str], size: int, bold: bool = False):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.clear()
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = str(line)
        para.font.size = Pt(size)
        para.font.bold = bold
    return slide
