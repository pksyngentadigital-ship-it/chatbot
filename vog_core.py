"""
Voice of Grower — framework-agnostic core.

Everything in this module is plain Python: no Streamlit imports, no UI calls.
It holds the constants, retrieval/ranking/trend logic, the Excel ingestion
parser, and the PPTX/Excel builders that power the chatbot. The Streamlit
app (app.py) is a thin UI wrapper around this module; a future FastAPI
backend can import the exact same module and get identical behavior.

Two-phase chat API (mirrors what any UI needs — get a prompt, stream the
LLM, then finalize with charts/downloads):

    state = process_chat_query(user_query, pinecone_api_key, groq_api_key)
    if state["kind"] in ("blocked", "no_key", "ranking", "trend", "no_data"):
        # state["reply"] is ready to display; state.get("downloads") may hold
        # ready CSV/Excel/PPTX bytes and state.get("chart") the chart data.
        ...
    elif state["kind"] == "normal":
        # Stream state["system_prompt"] / state["user_prompt"] through Groq
        # yourself (so each UI can render tokens its own way), then:
        result = finalize_normal_response(state, full_response_text)
        # result["final_reply"], result["chart"], result["downloads"]
"""

import hashlib
import json
import re
import os
from io import BytesIO
from collections import Counter

import pandas as pd
from pinecone import Pinecone
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ==========================================
# CONFIGURATION & CONSTANTS
# ==========================================
PINECONE_INDEX_NAME = "chatbot"
EMBEDDING_DIMENSION = 384
GROQ_MODEL = "openai/gpt-oss-20b"

# Fallback timeframe label used only when a query genuinely mentions no
# date at all — reads naturally in phrases like "Complaints of {label}:"
# and "Sentiments of {label}:". Never say the confusing literal "the
# requested period" back to a user who did specify a timeframe but whose
# phrasing the date parser didn't recognize.
DEFAULT_TIMEFRAME_LABEL = "all available feedback"

# Pinecone caps total metadata at 40KB per vector. A single Excel cell can
# hold 32k characters, and the old payload stored the bullet twice (once
# raw, once inside a longer "text" sentence), so one long cell could reject
# the whole 50-vector upsert batch it travelled in.
MAX_VALUE_CHARS = 8000

MONTH_TYPO_FIX = {
    "Feburary": "February", "Febuary": "February",
    "Septembar": "September", "Septmber": "September",
    "Octobar": "October",    "Novembar": "November",
    "Decembar": "December",  "Januray": "January",
    "Janaury": "January",    "Marck": "March"
}

MONTH_MAP = {
    "january": "January", "february": "February", "march": "March",
    "april": "April",     "may": "May",            "june": "June",
    "july": "July",       "august": "August",      "september": "September",
    "october": "October", "november": "November",  "december": "December",
    "jan": "January",     "feb": "February",       "mar": "March",
    "apr": "April",       "jun": "June",           "jul": "July",
    "aug": "August",      "sep": "September",      "oct": "October",
    "nov": "November",    "dec": "December"
}

MONTH_ORDER = {
    "January": 1, "February": 2, "March": 3, "April": 4, "May": 5, "June": 6,
    "July": 7, "August": 8, "September": 9, "October": 10, "November": 11,
    "December": 12
}

CATEGORY_NORMALIZE = {
    "product queries":              "Product Queries",
    "problem/advisory":             "Problem/Advisory",
    "problem advisory":             "Problem/Advisory",
    "positive feedback":            "Positive Feedback",
    "complaint/negative feedback":  "Complaint/Negative Feedback",
    "complaint negative feedback":  "Complaint/Negative Feedback",
    "complaints":                   "Complaint/Negative Feedback",
    "negative feedback":            "Complaint/Negative Feedback",
    "suggestion":                   "Suggestions",
    "suggestions":                  "Suggestions",
    "others":                       "Others",
    "other":                        "Others"
}

POSITIVE_CATEGORIES = {"Positive Feedback"}
# Problem/Advisory rows are grower-reported issues/concerns just like
# Complaint/Negative Feedback — both feed the "complaint" / "root cause" /
# "top concerns" style queries.
NEGATIVE_CATEGORIES = {"Complaint/Negative Feedback", "Problem/Advisory"}
SUGGESTION_CATEGORY  = "Suggestions"
PRODUCT_QUERY_CATEGORY = "Product Queries"
SALES_KEYWORDS = [
    "price", "pricing", "purchase", "buying", "unavailability",
    "unavailable", "sales-related", "cost of", "availability", "discount"
]

EMPTY_VALUES = {
    'nan', 'none', '', 'null', '-', 'n/a', 'na',
    'not filled', 'not available', 'no data', '0', 'tbd', 'pending'
}

# Known products — matched first (fast path). Extend freely.
PRODUCT_LIST = [
    # Verified against the real Naya Savera / Syngenta product catalog
    # (nayasavera.online) — both branded products and generic
    # active-ingredient names sold there.
    "cropwise", "quantis", "isabion", "isabion gold", "allymax", "axial",
    "walter", "walter super", "solubor", "amistar", "amistar top",
    "incipio", "simodis", "solvigo", "rifit", "logran", "cruiser", "enrich",
    "virtako", "proclaim", "thiovit", "thiovit jet", "thiovet",
    "pendimethalin", "polytrin", "polytrin c", "chlorpyrifos", "glyphosate",
    "tilt", "actara", "alika", "ridomil", "score", "folicur", "miraculan",
    "dual gold", "naya potash", "naya npk", "naya sop",
    "naya sulphate of potash", "naya s urea", "naya zinc plus", "promix",
    "promix npk", "karate", "dynasty", "dynasty cst", "buprofezin",
    "topas", "orondis", "orondis opti", "miravis", "miravis duo",
    "primextra", "primextra gold", "voliam flexi", "vibrance",
    "vibrance premium", "vibrance duo", "pyriproxyfin", "polo", "plenum",
    "revus", "revus start", "copper oxychloride", "metribuzin", "match",
    "gengwei", "bromoxynil", "dumei", "curacron", "bifenthrin", "ampligo",
    "acephate", "dragon", "elestal", "elestal neo"
]

# Known crops — matched first (fast path) for crop-wise analysis / filtering.
CROP_LIST = [
    "wheat", "rice", "paddy", "cotton", "maize", "corn", "sugarcane",
    "soybean", "soyabean", "groundnut", "mustard", "canola", "potato",
    "tomato", "onion", "chilli", "chili", "chickpea", "bengal gram",
    "pea", "peas", "banana", "grape", "grapes", "sunflower", "okra",
    "lady finger", "ladyfinger", "barley", "jowar", "bajra", "cabbage",
    "cauliflower", "brinjal", "cucumber", "watermelon", "muskmelon",
    "melon", "mango", "citrus", "orange", "apple", "turmeric", "ginger",
    "garlic", "sesame", "castor", "tobacco", "papaya", "guava",
    "pomegranate", "carrot", "radish", "spinach", "cumin", "coriander",
    "fenugreek",
    # NOTE: bare "gram" and "rose" are deliberately NOT in this list.
    # "gram" is a unit of mass and appears in nearly every dosage note
    # ("apply 50 gram per acre"), which made Gram a top-ranked crop;
    # "rose" is far more often the verb ("prices rose sharply"). The real
    # crops are reachable via "bengal gram"/"chickpea" and the guarded
    # pattern in CROP_GUARDED below.
]

# Crops whose bare name is ambiguous, matched only in unmistakably
# agronomic phrasing. Maps a regex to the canonical crop label.
CROP_GUARDED = {
    r'\b(?:bengal\s+gram|green\s+gram|black\s+gram|horse\s+gram)\b': "Chickpea",
    r'\brose\s+(?:crop|plants?|cultivation|growers?|farms?)\b': "Rose",
    r'\b(?:cut\s+)?roses\b': "Rose",
}

# Synonyms collapse to one canonical label so a ranking counts each crop
# once. Without this, rice(2)+paddy(2) lost the top slot to cotton(2).
CROP_CANONICAL = {
    "paddy": "Rice", "corn": "Maize", "soyabean": "Soybean",
    "chili": "Chilli", "grapes": "Grape", "peas": "Pea",
    "ladyfinger": "Okra", "lady finger": "Okra", "roses": "Rose",
    "chickpea": "Chickpea", "bengal gram": "Chickpea",
}

# Diseases, pests, and agronomic problems that must NEVER be reported back to
# the user as if they were products — the raw feedback text often names them
# in the same sentence as real products (e.g. "Septoria control with Score"),
# and the capitalized-phrase heuristic in extract_product_mentions would
# otherwise tag them as brand names.
DISEASE_PEST_TERMS = [
    "septoria", "blast", "blight", "rust", "wilt", "rot", "mildew",
    "powdery mildew", "downy mildew", "virus", "mosaic", "mosaic virus",
    "aphid", "aphids", "borer", "borers", "caterpillar", "caterpillars",
    "armyworm", "armyworms", "fall armyworm", "faw", "termite", "termites",
    "rodent", "rodents", "whitefly", "whiteflies", "thrips", "mite", "mites",
    "nematode", "nematodes", "weed", "weeds", "fungus", "fungal", "bacterial",
    "bacteria", "larvae", "larva", "infestation", "infestations", "scab",
    "canker", "leaf spot", "leaf curl", "yellowing", "stunting", "wilting",
    "disease", "diseases", "pest", "pests",
    # Pests that actually appear in this dataset. "Jassid" was ranking as
    # the single most-mentioned "product" in production purely because it
    # was missing from this list.
    "jassid", "jassids", "mealy", "mealybug", "mealybugs", "mealy bug",
    "bollworm", "bollworms", "pink bollworm", "american bollworm",
    "stem borer", "shoot borer", "fruit borer", "leaf miner", "leafminer",
    "grasshopper", "grasshoppers", "locust", "locusts", "cutworm",
    "cutworms", "semilooper", "semiloopers", "looper", "loopers",
    "anthracnose", "gummy stem blight", "damping off", "root rot",
    "collar rot", "sheath blight", "smut", "ergot", "tikka", "alternaria",
    "fusarium", "phytophthora", "sclerotinia", "botrytis",
    "nutrient deficiency", "deficiency", "chlorosis", "lodging",
]
# Flattened to individual words so multi-word phrases (e.g. "Leaf Curl Virus")
# are still caught even though only part of the phrase matches a listed term.
DISEASE_PEST_WORDS = {w for term in DISEASE_PEST_TERMS for w in term.split()}

# Extra business-domain vocabulary the strict topic guardrail must recognize
# (executive/sales/marketing/digital/advanced-analytics use cases, plus the
# various output-format requests such as tables, charts, and exports).
BUSINESS_KEYWORDS = [
    "suggestion", "suggestions", "recommend", "recommendation",
    "recommendations", "improvement", "improvements", "expectation",
    "expectations", "root", "cause", "causes", "trend", "trends",
    "trending", "rank", "ranking", "ranked", "top", "table", "excel",
    "export", "download", "ppt", "powerpoint", "presentation", "slide",
    "slides", "chart", "graph", "visualize", "visualise", "visualization",
    "visualisation", "dashboard", "executive", "summary", "summarize",
    "sales", "marketing", "digital", "campaign", "campaigns",
    "misconception", "misconceptions", "awareness", "yoy", "quarter",
    "quarterly", "insight", "insights", "strategic", "priority",
    "priorities", "forecast", "predict", "prediction", "region", "regions",
    "crop", "crops", "emerging", "satisfaction", "frequency", "frequent",
    "common", "pattern", "patterns", "hidden", "platform", "platforms",
    "online", "support", "experience", "customer", "recommended", "most",
    "highest", "significant", "significantly", "increased", "improve",
    "business", "monthly", "yearly", "annual", "annually", "breakdown",
    "revenue", "growth", "performance", "analytics", "kpi", "kpis"
]

# Generic words that should never be mistaken for a product name during
# dynamic (fallback) product detection.
PRODUCT_STOPWORDS = {
    "sentiment", "sentiments", "feedback", "feedbacks", "product", "products",
    "syngenta", "app", "price", "unavailability", "complaint", "complaints",
    "positive", "negative", "overview", "overall", "general", "summary",
    "both", "analysis", "grower", "growers", "farmer", "farmers", "people",
    "advisory", "week", "weeks",
    "list", "listed", "listing", "bullet", "bullets", "compare", "comparison",
    "versus", "give", "show", "tell", "what", "are", "about", "the", "for",
    "and", "of", "in", "on", "me", "please", "this", "that", "month",
    "months", "year", "years", "data", "record", "records", "issue",
    "issues", "problem", "problems", "concern", "concerns", "appreciation",
    "praise", "favorable", "satisfied", "first", "second", "third",
    "fourth", "fifth", "point", "points", "wise", "chatbot", "yield",
    # common connector / filler words that must never be treated as a
    # product name during dynamic detection
    "out", "down", "up", "into", "onto", "with", "from", "than", "then",
    "just", "only", "also", "very", "much", "many", "more", "most", "some",
    "such", "need", "needs", "want", "wants", "know", "get", "gets", "got",
    "can", "could", "would", "should", "will", "shall", "may", "might",
    "not", "no", "yes", "okay", "ok", "thanks", "thank", "you", "your",
    "our", "their", "his", "her", "its", "all", "any", "each", "every",
    "who", "whom", "which", "when", "where", "why", "how", "does", "did",
    "has", "have", "had", "was", "were", "been", "being", "here", "there",
    "these", "those", "over", "under", "again", "still", "yet", "now",
    "provide", "write", "respond", "answer", "reply", "query", "ask",
    "asking", "kindly", "regarding", "specific", "particular", "quick",
    "quickly", "brief", "detail", "details", "info", "information", "one",
    "two", "three", "four", "five", "recent", "latest", "last", "current",
    "previous", "next", "prior", "season", "seasons", "compared", "across",
    "during", "within", "improved", "day", "days", "reported", "report",
    "reports", "generated", "received", "found", "total", "number", "numbers",
    # "other"/"others" show up constantly in ordinary feedback prose (e.g.
    # "no other issues", "compared to other products") — without this, the
    # dynamic product-probe fallback mistakes the word itself for a product
    # name whenever a query like "what other insights..." reaches it.
    "other", "others", "another", "anything", "something", "else",
    # Phase 4 intent-keyword-expansion words (complaint/positive/suggestion/
    # sentiment phrasing added below in process_chat_query) are common
    # English words that show up verbatim in real feedback text — same
    # false-positive risk "other"/"pricing" already exposed live. "delayed"
    # was confirmed live to trigger this exact bug.
    "dissatisfied", "unhappy", "frustration", "frustrated", "shortage",
    "shortages", "delay", "delays", "delayed", "defect", "defects",
    "faulty", "damaged", "working", "poor", "low", "bad", "quality",
    "disappointed", "difficulty", "difficulties", "trouble", "troubles",
    "grievance", "grievances", "happy", "pleased", "impressed", "love",
    "loved", "loving", "great", "worked", "works", "well", "effective",
    "satisfaction", "delighted", "thrilled", "wish", "wishes", "hope",
    "hoping", "request", "requests", "requested", "see", "include",
    "feature", "enhancement", "enhancements", "think", "thoughts",
    "opinion", "opinions", "views", "perception", "perceptions",
    "reaction", "reactions", "impression", "impressions", "experience",
    "experiences",
    # "Kaho" is a Syngenta podcast/campaign name, not a product — it was
    # mistakenly included in an earlier PRODUCT_LIST pass and genuinely
    # appears in real feedback text (podcast mentions), so it must be
    # excluded here too or the dynamic product-probe fallback would just
    # re-confirm it as a "product" anyway.
    "kaho",
} | set(MONTH_MAP.keys()) | set(BUSINESS_KEYWORDS) | set(DISEASE_PEST_TERMS) | set(SALES_KEYWORDS) | set(CROP_LIST) | {
    # A crop is never a product. "Anthracnose in Chilli crop" was tagging
    # Chilli as a brand, which then competed with real products in the
    # product ranking.
    w for crop in CROP_LIST for w in crop.split()
}

ALLOWED_GUARDRAIL_KEYWORDS = set([
    "sentiment", "sentiments", "feedback", "feedbacks", "product", "products",
    "syngenta", "cropwise", "app", "price", "unavailability", "january",
    "february", "march", "april", "may", "june", "july", "august", "september",
    "october", "november", "december", "jan", "feb", "mar", "apr", "jun", "jul",
    "aug", "sep", "oct", "nov", "dec", "2023", "2024", "2025", "2026", "2027",
    "complaint", "complaints", "positive", "negative", "grower", "growers",
    "advisory", "week", "weeks", "1st", "2nd", "3rd", "4th", "5th", "first",
    "second", "third", "fourth", "fifth", "issues", "concerns", "problems",
    "appreciation", "praise", "compare", "comparison", "versus", "list"
]) | set(PRODUCT_LIST) | set(CROP_LIST) | set(BUSINESS_KEYWORDS)

SUGGESTED_PROMPTS = {
    "📊 Executive": [
        "What are the top 10 grower concerns reported during the last 30 days?",
        "Summarize the major grower insights for the last quarter in one page.",
        "What are the emerging trends compared to the previous year?",
    ],
    "🌾 Crop-wise": [
        "What are the top issues reported in rice during 2026?",
        "Compare grower feedback for wheat between 2025 and 2026.",
        "Show the sentiment analysis for cotton growers.",
    ],
    "🏷️ Product-wise": [
        "Which Syngenta products received the highest positive feedback?",
        "List the products with the highest complaint frequency.",
        "Compare customer sentiment for Tilt and Isabion.",
    ],
    "🐛 Complaints": [
        "What are the top five complaints received in the last six months?",
        "Identify the root causes of the most common complaints.",
        "Show the monthly complaint trend for the past three years.",
    ],
    "🌻 Sentiment": [
        "Show overall grower sentiment by month.",
        "Compare positive, neutral, and negative sentiment across crops.",
        "Display sentiment trends in chart and table format.",
    ],
    "💡 Suggestions": [
        "What suggestions have growers provided for rice cultivation?",
        "What are the most common product improvement recommendations?",
        "Summarize grower expectations from Syngenta.",
    ],
    "💰 Sales": [
        "Which products receive the highest number of price inquiries?",
        "Which crops generate the highest sales-related questions?",
    ],
    "📣 Marketing & Digital": [
        "What misconceptions do growers commonly have about our products?",
        "What are the most common complaints regarding digital platforms or systems?",
    ],
    "🔮 Advanced & Rankings": [
        "Which crop generated the highest number of complaints?",
        "Generate the Top 10 business insights from the last three years.",
        "Based on all historical data, recommend the top five strategic actions Syngenta should prioritize for the next season.",
    ],
    "📄 Output formats": [
        "Show grower suggestions for rice in a table.",
        "Generate an executive summary of complaints for 2026.",
        "Create a PowerPoint-ready summary of positive feedback for Isabion.",
    ],
}

# A short, curated subset of SUGGESTED_PROMPTS shown on the FastAPI chat
# UI's welcome screen — one high-value prompt per major use case, so first-
# time users see a handful of clear options instead of the full category
# grid. (The Streamlit app still uses the full SUGGESTED_PROMPTS above.)
SUGGESTED_PROMPTS_QUICK = [
    "What are the top 10 grower concerns reported during the last 30 days?",
    "Which Syngenta products received the highest positive feedback?",
    "Show overall grower sentiment by month.",
    "What are the most common product improvement recommendations?",
    "Which crop generated the highest number of complaints?",
    "Summarize the major grower insights for the last quarter in one page.",
]

# ==========================================
# UTILITIES
# ==========================================

def find_category_column(df_columns):
    """Prefer an exact 'category' header. Falling back to the first header
    merely CONTAINING 'categ' picked 'Sub Category' over 'Category' and
    'Categorization Notes' over the real column — and the wrong column then
    drove the entire Layout A branch, tagging every record with a value that
    normalizes to nothing and therefore lands as sentiment='neutral'."""
    exact_targets = {"category", "case category", "casecategory"}
    for col in df_columns:
        if re.sub(r'\s+', ' ', str(col)).strip().lower() in exact_targets:
            return col
    for col in df_columns:
        if 'categ' in re.sub(r'\s+', '', str(col)).lower():
            return col
    return None


def infer_year_for_sheet(sheet_name: str, all_sheet_names: list) -> str | None:
    """ Sheets carrying an explicit 4-digit year are trusted directly. Undated legacy sheets ("Jan till June") are dated from their NEAREST dated neighbour by position — the previous implementation took the global min of all later sheets / max of all earlier ones, which for ['Legacy', 'VOG 2026', 'VOG 2025'] produced 2024 instead of 2025, and could silently assign a legacy sheet the same year as a real dated sheet, double-counting every record in it. Returns None rather than guessing when the inferred year would collide with an explicitly-dated sheet. """
    direct = re.search(r'\b(20\d{2})\b', sheet_name.strip())
    if direct:
        return direct.group(1)

    dated = [
        (i, int(m.group(1)))
        for i, name in enumerate(all_sheet_names)
        for m in [re.search(r'\b(20\d{2})\b', str(name).strip())] if m
    ]
    if not dated:
        return None

    my_index = None
    for i, name in enumerate(all_sheet_names):
        if str(name).strip() == sheet_name.strip():
            my_index = i
            break
    if my_index is None:
        return None

    # Nearest dated sheet by positional distance, not global min/max.
    nearest_i, nearest_year = min(dated, key=lambda t: (abs(t[0] - my_index), t[0]))
    guess = nearest_year - 1 if nearest_i > my_index else nearest_year + 1

    claimed = {year for _, year in dated}
    if guess in claimed:
        # Refuse rather than duplicate an explicitly-dated sheet's year.
        return None
    return str(guess)


def normalize_category(raw_val):
    """Returns the canonical category name, or None when the value does not
    map to a known category. Previously an unknown spelling was passed
    through verbatim and then silently tagged sentiment='neutral' — so
    'Complaint / Negative Feedback' (spaces around the slash) made an entire
    sheet of complaints invisible to every complaint query. Callers now get
    None and are expected to record it as a skipped/unmapped value."""
    if raw_val is None:
        return None
    cleaned = re.sub(r'\s*/\s*', '/', re.sub(r'\s+', ' ', str(raw_val)).strip().lower())
    if not cleaned or cleaned in EMPTY_VALUES:
        return None
    if cleaned in CATEGORY_NORMALIZE:
        return CATEGORY_NORMALIZE[cleaned]
    # Tolerate a trailing plural ("Positive Feedbacks").
    if cleaned.endswith("s") and cleaned[:-1] in CATEGORY_NORMALIZE:
        return CATEGORY_NORMALIZE[cleaned[:-1]]
    return None


def is_empty_cell(value: str) -> bool:
    return value.strip().lower() in EMPTY_VALUES or value.strip() == ""


def split_bullets(cell_text: str) -> list[str]:
    """Split a cell into individual feedback points. Splits on newlines AND
    on inline bullet glyphs / sentence-ending semicolons — a cell authored
    as "• A • B • C" on one line used to become a single record with one
    blended set of crop/product tags."""
    parts = re.split(r'[\n\r]+|\s+[•●·]\s*|(?<=[.;])\s+(?=[A-Z0-9])', cell_text)
    bullets = []
    for line in parts:
        if line is None:
            continue
        clean = re.sub(r'^[\s•●·\-–—]+', '', line).strip()
        if clean and clean.lower() not in EMPTY_VALUES and len(clean) >= 2:
            bullets.append(clean)
    return bullets


_MONTH_RE = re.compile(
    r'\b(january|february|march|april|may|june|july|august'
    r'|september|october|november|december'
    r'|jan|feb|mar|apr|jun|jul|aug|sep|oct|nov|dec)\b',
    re.IGNORECASE,
)


def extract_month_from_col(col: str) -> str:
    """Word-boundary month match. Without \\b the alternation matched inside
    ordinary words — 'Remarks' -> March, 'Summary' -> March,
    'Decrease' -> December, 'Separate' -> September — which filed real rows
    under a month nobody wrote."""
    fixed = str(col)
    for typo, correct in MONTH_TYPO_FIX.items():
        fixed = re.sub(r'\b' + typo + r'\b', correct, fixed, flags=re.IGNORECASE)
    match = _MONTH_RE.search(fixed)
    if not match:
        return "Unknown"
    return MONTH_MAP.get(match.group(1).lower(), match.group(1).capitalize())


# A single record written at the end of ingestion holding the dataset's
# real extents. Reading these from a sampled query was unsound: a zero
# vector carries no similarity signal, so Pinecone returns arbitrary
# records, and taking max(year) over 10 of them could resolve "March" to
# the wrong year entirely — non-deterministically, while the header stated
# that wrong year with full confidence.
INDEX_STATS_ID = "vog_index_stats"


def _read_index_stats(index) -> dict | None:
    try:
        res = index.fetch(ids=[INDEX_STATS_ID])
        vectors = getattr(res, "vectors", None) or (res.get("vectors") if isinstance(res, dict) else None) or {}
        entry = vectors.get(INDEX_STATS_ID)
        if not entry:
            return None
        md = entry.get("metadata") if isinstance(entry, dict) else getattr(entry, "metadata", None)
        return md or None
    except Exception:
        return None


def get_latest_year_from_index(index) -> str:
    stats = _read_index_stats(index)
    if stats and str(stats.get("max_year", "")).isdigit():
        return str(stats["max_year"])

    # Fallback for an index ingested before stats were written. Sample far
    # more widely than the old top_k=10, which was small enough to miss the
    # most recent year outright.
    try:
        dummy_vector = [0.0] * EMBEDDING_DIMENSION
        results = index.query(vector=dummy_vector, top_k=1000, include_metadata=True)
        years = [
            int(y) for m in results.get("matches", [])
            for y in [m.get("metadata", {}).get("year")]
            if y and str(y).isdigit()
        ]
        if years:
            return str(max(years))
    except Exception:
        pass
    return "2026"


def get_max_week_label(index, month, year) -> str | None:
    """Find the actual latest week label (e.g. '5th Week') stored in the dataset for the given month/year, so 'last week' / 'latest week' queries resolve to a real week instead of a fixed guess."""
    filter_conditions = {}
    if month:
        filter_conditions["month"] = {"$eq": month}
    if year:
        filter_conditions["year"] = {"$eq": year}
    try:
        dummy_vector = [0.0] * EMBEDDING_DIMENSION
        results = index.query(
            vector=dummy_vector, top_k=200, include_metadata=True,
            filter=filter_conditions if filter_conditions else None
        )
        weeks = [m.get("metadata", {}).get("week", "") for m in results.get("matches", [])]
        weeks = [w for w in weeks if w]
        if not weeks:
            return None

        def week_num(w):
            match = re.search(r'(\d+)', w)
            return int(match.group(1)) if match else -1

        return max(set(weeks), key=week_num)
    except Exception:
        return None


MONTH_ORDER_INV = {v: k for k, v in MONTH_ORDER.items()}


def get_latest_month_year_from_index(index) -> tuple[str, str] | None:
    """ Find the most recent (month, year) pair actually present in the ingested data — the anchor every relative-date phrase ("last month", "last quarter") resolves against. Anchored to the data, not the real calendar date, for the same reason get_latest_year_from_index is: the dataset may lag behind today, so a real-calendar "last month" could return empty even when recent data exists. """
    stats = _read_index_stats(index)
    if stats:
        month, year = stats.get("max_month"), stats.get("max_year")
        if month in MONTH_ORDER and str(year).isdigit():
            return month, str(year)

    try:
        dummy_vector = [0.0] * EMBEDDING_DIMENSION
        results = index.query(vector=dummy_vector, top_k=1000, include_metadata=True)
        pairs = []
        for m in results.get("matches", []):
            md = m.get("metadata", {})
            month, year = md.get("month"), md.get("year")
            if month in MONTH_ORDER and year and str(year).isdigit():
                pairs.append((int(year), month))
        if not pairs:
            return None
        latest_year, latest_month = max(pairs, key=lambda p: (p[0], MONTH_ORDER[p[1]]))
        return latest_month, str(latest_year)
    except Exception:
        return None


def _month_idx(month: str, year: str) -> int:
    """Absolute month index (months since year 0) — makes month arithmetic a plain integer add/subtract instead of manual year-rollover juggling."""
    return int(year) * 12 + (MONTH_ORDER[month] - 1)


def _idx_to_month_year(idx: int) -> tuple[str, str]:
    year, month0 = divmod(idx, 12)
    return MONTH_ORDER_INV[month0 + 1], str(year)


def detect_relative_window(query_lower: str) -> dict | None:
    """ Detects relative time-window phrases that span MULTIPLE months — "last 30 days", "last quarter", "last 3 months", "since March" — which the exact single month/year/week matchers can't express at all today. Returns a small descriptor dict for resolve_relative_window(), or None. Note on granularity: the ingested data only carries month/week tags, no exact calendar date — so "last N days" and "last N weeks" are necessarily approximated to the nearest whole month count (round(N/30) and round(N/4) respectively). That's a real, documented limitation of the data model, not a rounding bug. """
    m = re.search(r'\b(?:last|past)\s+(\d+)\s+days?\b', query_lower)
    if m:
        months_back = max(1, round(int(m.group(1)) / 30))
        return {"kind": "last_n_months", "n": months_back}

    m = re.search(r'\b(?:last|past)\s+(\d+)\s+weeks?\b', query_lower)
    if m:
        months_back = max(1, round(int(m.group(1)) / 4))
        return {"kind": "last_n_months", "n": months_back}

    m = re.search(r'\b(?:last|past|previous)\s+(\d+)\s+months?\b', query_lower)
    if m:
        return {"kind": "last_n_months", "n": int(m.group(1))}

    if re.search(r'\blast\s+month\b', query_lower):
        return {"kind": "last_n_months", "n": 1}

    if re.search(r'\bthis\s+month\b|\bcurrent\s+month\b', query_lower):
        return {"kind": "this_month"}

    if re.search(r'\b(?:last|past|previous)\s+quarter\b', query_lower):
        return {"kind": "last_quarter"}

    if re.search(r'\bthis\s+quarter\b|\bcurrent\s+quarter\b', query_lower):
        return {"kind": "this_quarter"}

    m = re.search(
        r'\bsince\s+(january|february|march|april|may|june|july|august'
        r'|september|october|november|december'
        r'|jan|feb|mar|apr|jun|jul|aug|sep|oct|nov|dec)\b',
        query_lower
    )
    if m:
        return {"kind": "since_month", "month": MONTH_MAP.get(m.group(1), m.group(1).capitalize())}

    return None


def _interleave_by_recency(per_month_bullets: list[list[str]]) -> list[str]:
    """ Merge per-month bullet lists (given oldest-month-first) into one list that stays representative after downstream truncation. Round-robins one bullet from each month starting with the NEWEST, so a "last N months" answer covers the whole window with a recency bias, instead of being filled entirely by whichever month happens to come first. Deduplicates case-insensitively. """
    merged, seen = [], set()
    newest_first = list(reversed(per_month_bullets))
    for i in range(max((len(b) for b in newest_first), default=0)):
        for month_bullets in newest_first:
            if i < len(month_bullets):
                b = month_bullets[i]
                key = b.strip().lower()
                if key not in seen:
                    seen.add(key)
                    merged.append(b)
    return merged


def resolve_relative_window(window: dict, index) -> tuple[list[tuple[str, str]], str] | None:
    """ Resolves a detect_relative_window() descriptor into an ordered list of (month, year) pairs to query and merge, plus a human-readable label for the response header. Anchored to the latest month/year actually in the data (see get_latest_month_year_from_index) — never the real calendar date. Returns None if the index has no dated records to anchor against. """
    latest = get_latest_month_year_from_index(index)
    if not latest:
        return None
    latest_month, latest_year = latest
    latest_idx = _month_idx(latest_month, latest_year)

    kind = window["kind"]
    if kind == "last_n_months":
        n = max(1, window["n"])
        idxs = list(range(latest_idx - n + 1, latest_idx + 1))
        label = f"the last {n} month{'s' if n != 1 else ''}"
    elif kind == "this_month":
        idxs = [latest_idx]
        label = f"{latest_month} {latest_year}"
    elif kind == "last_quarter":
        cur_q_start = (latest_idx // 3) * 3
        q_start = cur_q_start - 3
        idxs = list(range(q_start, q_start + 3))
        label = "last quarter"
    elif kind == "this_quarter":
        cur_q_start = (latest_idx // 3) * 3
        idxs = list(range(cur_q_start, latest_idx + 1))
        label = "this quarter"
    elif kind == "since_month":
        target_idx = _month_idx(window["month"], latest_year)
        if target_idx > latest_idx:
            target_idx -= 12  # that month hasn't happened yet this year — use last year's
        idxs = list(range(target_idx, latest_idx + 1))
        label = f"since {window['month']}"
    else:
        return None

    months = [_idx_to_month_year(i) for i in idxs]
    return months, label


def query_pinecone_for_timeframe(index, query_vector, month, year, week, query_intent="sentiment", top_k=100, category_filter=None):
    filter_conditions = {}
    if month:
        filter_conditions["month"] = {"$eq": month}
    if year:
        filter_conditions["year"]  = {"$eq": year}

    # Week as a DATABASE filter, not a Python post-filter. Applying it after
    # the top_k cut meant the similarity ranking selected 100 records for the
    # whole month first, and the week filter then discarded most of what
    # came back — returning "no data" for weeks that genuinely had data,
    # silently and in proportion to how large the month was.
    week_num = _week_number(week) if week else None
    if week_num is not None:
        filter_conditions["week_num"] = {"$eq": week_num}

    # ── CATEGORY FILTER (e.g. "Suggestions") TAKES PRIORITY OVER SENTIMENT ──
    if category_filter:
        filter_conditions["category"] = {"$eq": category_filter}
    # ── SENTIMENT FILTER AT DATABASE LEVEL ──
    elif query_intent == "positive":
        filter_conditions["sentiment"] = {"$eq": "positive"}
    elif query_intent == "complaint":
        filter_conditions["sentiment"] = {"$eq": "negative"}

    metadata_filter = filter_conditions if filter_conditions else None

    results = index.query(
        vector=query_vector,
        top_k=top_k,
        include_metadata=True,
        filter=metadata_filter
    )
    matches = results.get("matches", [])

    # Records written before week_num existed have no such field and are
    # excluded by the DB filter above, so fall back to the legacy text
    # match when the filtered query comes back empty.
    if week_num is not None and not matches:
        legacy_filter = {k: v for k, v in filter_conditions.items() if k != "week_num"}
        results = index.query(
            vector=query_vector, top_k=top_k, include_metadata=True,
            filter=legacy_filter or None,
        )
        matches = results.get("matches", [])

    ORDINAL_MAP = {
        "first": "1st", "second": "2nd",
        "third": "3rd", "fourth": "4th", "fifth": "5th"
    }

    positive_bullets = []
    negative_bullets = []
    neutral_bullets  = []

    # Track exact text already added (case-insensitive) so the same feedback
    # point is never sent to the model twice — this covers duplicate vectors
    # from re-ingesting the same sheet, or the same line matching more than
    # one week column.
    seen_positive = set()
    seen_negative = set()
    seen_neutral  = set()

    for m in matches:
        md        = m.get("metadata", {})
        sentiment = md.get("sentiment", "neutral")
        value     = md.get("value", "").strip()
        w         = md.get("week", "")
        category  = md.get("category", "")

        if not value or value.lower() in EMPTY_VALUES:
            continue

        if week:
            dw = week.lower()
            for word, num in ORDINAL_MAP.items():
                dw = dw.replace(word, num)
            if dw not in w.lower():
                continue

        entry = f"{category}: {value}"
        dedupe_key = entry.strip().lower()

        if sentiment == "positive":
            if dedupe_key not in seen_positive:
                seen_positive.add(dedupe_key)
                positive_bullets.append(entry)
        elif sentiment == "negative":
            if dedupe_key not in seen_negative:
                seen_negative.add(dedupe_key)
                negative_bullets.append(entry)
        else:
            if dedupe_key not in seen_neutral:
                seen_neutral.add(dedupe_key)
                neutral_bullets.append(entry)

    return positive_bullets, negative_bullets, neutral_bullets


def extract_all_months(query_lower: str) -> list[str]:
    """Return every distinct month mentioned in the query, in the order first seen."""
    found = []
    for shortcut in sorted(MONTH_MAP.keys(), key=len, reverse=True):
        if re.search(r'\b' + re.escape(shortcut) + r'\b', query_lower):
            full = MONTH_MAP[shortcut]
            if full not in found:
                found.append(full)
    return found


def extract_all_years(query_lower: str) -> list[str]:
    """Return every distinct 4-digit year mentioned, in order first seen."""
    return list(dict.fromkeys(re.findall(r'\b(20\d{2})\b', query_lower)))


def extract_all_weeks(query_lower: str) -> list[str]:
    """Return every distinct 'Nth week' phrase mentioned, normalized."""
    ORDINAL_MAP = {
        "first": "1st", "second": "2nd",
        "third": "3rd", "fourth": "4th", "fifth": "5th"
    }
    raw_matches = re.findall(
        r'\b(1st|2nd|3rd|4th|5th|first|second|third|fourth|fifth)\s+week\b',
        query_lower
    )
    normalized = []
    for m in raw_matches:
        val = ORDINAL_MAP.get(m, m)
        if val not in normalized:
            normalized.append(val)
    return normalized


def detect_product_known(query_lower: str) -> str | None:
    """Fast path: match against the curated PRODUCT_LIST."""
    for product in PRODUCT_LIST:
        if re.search(r'\b' + re.escape(product) + r'\b', query_lower):
            return product
    return None


# A candidate must look like a brand rather than an ordinary word before it
# is worth probing at all, and it must then be corroborated by the curated
# `products` metadata tag rather than by appearing anywhere in free text.
MAX_DYNAMIC_CANDIDATES = 3
MAX_QUERY_WORDS_FOR_PROBE = 30
MIN_PRODUCT_TAG_HITS = 2


def _dynamic_candidates(query_lower: str, original_query: str) -> list[str]:
    """Candidate ordering for the dynamic probe: longest first (a brand name
    is usually the most distinctive token in the sentence), and capitalized-
    in-the-original tokens preferred, since real product names are written
    as proper nouns. Deliberately NOT query order — that made a leading verb
    like "provided" outrank the actual product name."""
    words = re.findall(r'\b[a-zA-Z]{4,}\b', query_lower)
    seen, cands = set(), []
    capitalized = {w.lower() for w in re.findall(r'\b[A-Z][a-zA-Z]{3,}\b', original_query)}
    for w in words:
        if w in PRODUCT_STOPWORDS or w in seen:
            continue
        seen.add(w)
        cands.append(w)
    cands.sort(key=lambda w: (w in capitalized, len(w)), reverse=True)
    return cands


def detect_product_dynamic(query_lower: str, index, pc, original_query: str = "") -> str | None:
    """ Fallback path for products NOT in PRODUCT_LIST. Confirms a candidate against the curated `products` metadata tag written at ingestion — NOT against raw feedback text. The old substring-in-free-text test confirmed almost any common English word, because in a corpus made entirely of grower feedback nearly every word appears somewhere; that produced a steady stream of false products ("other", "pricing", "delayed", "farmers", and — from the app's own suggested prompts — "provided" and "past"). Also bounded: at most MAX_DYNAMIC_CANDIDATES probes, skipped entirely for long queries, and all candidate embeddings requested in ONE batched call rather than one round trip per word. """
    if not query_lower.strip():
        return None
    # A long query is prose, not a product lookup — and probing it used to
    # mean one embedding call plus one vector query per word.
    if len(re.findall(r'\b\w+\b', query_lower)) > MAX_QUERY_WORDS_FOR_PROBE:
        return None

    candidates = _dynamic_candidates(query_lower, original_query or query_lower)[:MAX_DYNAMIC_CANDIDATES]
    if not candidates:
        return None

    try:
        embed_response = pc.inference.embed(
            model="llama-text-embed-v2",
            inputs=[f"{c} product feedback sentiment" for c in candidates],
            parameters={"input_type": "query", "dimension": EMBEDDING_DIMENSION}
        )
        vectors = [item.values for item in embed_response]
    except Exception:
        return None

    for cand, vector in zip(candidates, vectors):
        try:
            probe = index.query(vector=vector, top_k=50, include_metadata=True)
        except Exception:
            continue

        # Corroborate against the curated product tag, with a word-boundary
        # match and a minimum hit count — a genuine product name dominates
        # its own probe; a common word shows up diffusely or not at all.
        pattern = re.compile(r'(?:^|,)\s*' + re.escape(cand) + r'\s*(?:,|$)', re.IGNORECASE)
        hits = sum(
            1 for m in probe.get("matches", [])
            if pattern.search(str(m.get("metadata", {}).get("products", "")))
        )
        if hits >= MIN_PRODUCT_TAG_HITS:
            return cand
    return None


def _mentions(text: str, term: str) -> bool:
    """Word-boundary containment. Bare substring matching made 'rice' match
    'price' and 'gram' match 'program' — and it was inconsistent with
    extract_crops, which has always used \\b at ingestion time."""
    return bool(re.search(r'\b' + re.escape(term.lower()) + r'\b', text.lower()))


def filter_bullets_by_product(bullets: list[str], product: str) -> list[str]:
    """Keep only bullets that actually reference the requested product."""
    return [b for b in bullets if _mentions(b, product)]


def detect_crop(query_lower: str) -> str | None:
    """Fast path: match against the curated CROP_LIST (closed vocabulary, no dynamic probe needed)."""
    for crop in CROP_LIST:
        if re.search(r'\b' + re.escape(crop) + r'\b', query_lower):
            return crop
    return None


def filter_bullets_by_crop(bullets: list[str], crop: str) -> list[str]:
    """Keep only bullets that actually reference the requested crop, allowing
    for synonyms (a 'rice' query should also match bullets saying 'paddy')."""
    aliases = {crop.lower()} | {
        k for k, v in CROP_CANONICAL.items() if v.lower() == canonical_crop(crop).lower()
    } | {canonical_crop(crop).lower()}
    return [b for b in bullets if any(_mentions(b, a) for a in aliases)]


def canonical_crop(crop: str) -> str:
    """Collapse crop synonyms (paddy->Rice, corn->Maize) to one label."""
    return CROP_CANONICAL.get(crop.strip().lower(), crop.strip().title())


def extract_crops(text: str) -> list[str]:
    """Tag every known crop mentioned in a feedback bullet, for ingestion-time
    metadata. Synonyms are collapsed so each real crop is counted once."""
    text_lower = text.lower()
    found = []
    for crop in CROP_LIST:
        if re.search(r'\b' + re.escape(crop) + r'\b', text_lower):
            label = canonical_crop(crop)
            if label not in found:
                found.append(label)
    for pattern, label in CROP_GUARDED.items():
        if re.search(pattern, text_lower) and label not in found:
            found.append(label)
    return found


# Common transition/adjective words that show up capitalized purely because
# they start a sentence (e.g. "Excellent results with Isabion.", "However,
# growers..."). These previously polluted the product-ranking feature with
# entries like "Excellent", "Best", "However", "Add".
GENERIC_CAPITALIZED_STOPWORDS = {
    "however", "therefore", "additionally", "furthermore", "also",
    "moreover", "meanwhile", "otherwise", "instead", "besides", "thus",
    "hence", "accordingly", "consequently", "nonetheless", "nevertheless",
    "add", "added", "adding", "train", "training", "commerce", "delivery",
    "excellent", "good", "best", "great", "nice", "poor", "bad", "better",
    "worse", "worst", "outstanding", "exceptional", "effective",
    "ineffective", "satisfactory", "unsatisfactory", "quality", "results",
    "result", "perfect", "amazing", "wonderful", "fantastic", "impressive",
    "disappointing", "happy", "unhappy", "pleased", "displeased", "overall",
    "regarding", "concerning", "unfortunately", "fortunately", "currently",
    "recently", "generally", "specifically", "basically", "actually",
    "certainly", "definitely", "probably", "possibly", "apparently",
    "clearly", "obviously", "importantly", "essentially",
    # Agronomy-advice vocabulary that reads as a capitalized noun phrase in
    # this dataset ("Solution for jassid in cotton") but is never a brand.
    "solution", "solutions", "control", "management", "treatment",
    "dosage", "dose", "spray", "application", "recommendation", "advisory",
    "attack", "damage", "stage", "crop", "field", "farmer", "farmers",
    "grower", "growers", "market", "dealer", "retailer", "distributor",
}
PRODUCT_STOPWORDS |= GENERIC_CAPITALIZED_STOPWORDS


# Catalog entries that are also ordinary English words. Matching these
# case-insensitively tagged "did not match the description" as the product
# Match, and "NPS score dropped" as Score. They only count as products when
# capitalized in the source text.
AMBIGUOUS_PRODUCT_WORDS = {
    "match", "score", "enrich", "dragon", "tilt", "polo", "walter",
    "cruiser", "karate", "dynasty", "revus", "plenum",
}


def _canonical_product(product: str) -> str:
    return " ".join(
        w.upper() if w.lower() in ("sop", "cst", "npk") else w.capitalize()
        for w in product.split()
    )


def extract_product_mentions(text: str) -> list[str]:
    """ Tag likely product/brand names in a feedback bullet, for the ingestion-time metadata that deterministic ranking counts. Pass 1 matches the curated PRODUCT_LIST; pass 2 is a capitalized-phrase heuristic for brands not yet catalogued. Three corrections over the naive version: a catalog match that is a strict prefix of another match on the same bullet is dropped (so "Isabion Gold" no longer also credits "Isabion", which inflated base names and split variant counts); catalog entries that are ordinary English words must be capitalized in the source; and a capitalized phrase is rejected if ANY of its words is a stopword (the previous `all(...)` let "PRICE TOO HIGH" through as a brand because only "price" was listed). """
    text_lower = text.lower()
    seen = set()

    catalog_hits = []
    for product in PRODUCT_LIST:
        if not re.search(r'\b' + re.escape(product) + r'\b', text_lower):
            continue
        if product in AMBIGUOUS_PRODUCT_WORDS:
            # Require the capitalized form in the original text.
            if not re.search(r'\b' + re.escape(product.title()) + r'\b', text):
                continue
        catalog_hits.append(product)

    # Drop any catalog hit that is a strict prefix of a longer hit.
    filtered = [
        p for p in catalog_hits
        if not any(other != p and other.startswith(p + " ") for other in catalog_hits)
    ]

    out = []
    for product in filtered:
        canonical = _canonical_product(product)
        if canonical.lower() not in seen:
            seen.add(canonical.lower())
            out.append(canonical)

    candidates = re.findall(r'\b([A-Z][a-zA-Z]{2,}(?:\s+[A-Z][a-zA-Z]{2,}){0,2})\b', text)
    for cand in candidates:
        words = cand.split()
        lowered = [w.lower() for w in words]
        # ANY stopword disqualifies the phrase, not only all-of-them.
        if any(w in PRODUCT_STOPWORDS or w in MONTH_MAP for w in lowered):
            continue
        if any(w in DISEASE_PEST_WORDS for w in lowered):
            continue
        if any(w in GENERIC_CAPITALIZED_STOPWORDS for w in lowered):
            continue
        if cand.strip().lower() in ("syngenta",):
            continue
        # Multi-word ALL-CAPS is shouted feedback ("PRICE TOO HIGH"), not a brand.
        if len(words) > 1 and cand.isupper():
            continue
        key = cand.strip().lower()
        if key not in seen:
            seen.add(key)
            out.append(cand.strip())
    return out


def build_comparison_periods(all_months, all_years, all_weeks, index):
    """ Build a list of (label, month, year, week) tuples describing each period to compare. The dimension with 2+ distinct values becomes the axis of comparison; other dimensions are held fixed. No explicit "compare" keyword is required — mentioning two+ months/years/weeks is enough. """
    periods = []

    if len(all_years) >= 2:
        month = all_months[0] if all_months else None
        week  = all_weeks[0] if all_weeks else None
        for y in all_years:
            label = f"{(month + ' ') if month else ''}{y}"
            periods.append((label, month, y, week))

    elif len(all_months) >= 2:
        year = all_years[0] if all_years else get_latest_year_from_index(index)
        week = all_weeks[0] if all_weeks else None
        for m in all_months:
            label = f"{m} {year}"
            periods.append((label, m, year, week))

    elif len(all_weeks) >= 2:
        month = all_months[0] if all_months else None
        year  = all_years[0] if all_years else get_latest_year_from_index(index)
        for w in all_weeks:
            label = f"{w} week" + (f" of {month}" if month else "") + f" {year}"
            periods.append((label, month, year, w))

    return periods


def detect_aggregation_request(query_lower: str) -> str | None:
    """ Detects queries that want a deterministic, counted ranking ("which crop generated the highest number of complaints", "products with the highest complaint frequency") rather than free-form LLM summarization. Returns 'crop' or 'product' — the metadata field to rank by — or None. Kept separate from the normal retrieval flow because counting must be exact (computed in Python from real metadata tags), never left to the LLM to eyeball from a handful of retrieved bullets. Requires crop/product to be the explicit SUBJECT of the ranking (e.g. "which crop...", "products with the highest...") — a bare co-occurrence of a rank word and the word "product" isn't enough, since that also matches open-ended synthesis asks like "most common product improvement recommendations", which want an LLM to cluster free text, not a tag count. """
    synthesis_markers = [
        'recommend', 'suggestion', 'insight', 'improvement', 'expectation',
        'root cause', 'strategic action',
    ]
    if any(m in query_lower for m in synthesis_markers):
        return None

    crop_subject = bool(
        re.search(r'\bwhich\s+crops?\b', query_lower)
        or re.search(r'\bcrops?\s+(?:with|that|generated|received|has|have)\b', query_lower)
        or re.search(r'\btop\s+\d+\s+crops?\b', query_lower)
        or re.search(r'\bcrop[- ]wise\s+ranking\b', query_lower)
    )
    product_subject = bool(
        re.search(r'\bwhich\s+products?\b', query_lower)
        or re.search(r'\bproducts?\s+(?:with|that|generated|received|has|have)\b', query_lower)
        or re.search(r'\btop\s+\d+\s+products?\b', query_lower)
        or re.search(r'\bproduct[- ]wise\s+ranking\b', query_lower)
    )
    if not (crop_subject or product_subject):
        return None

    rank_phrases = [
        'highest number of', 'highest', 'most frequent', 'most complaints',
        'most common', 'ranking', 'rank ', 'frequency', 'greatest'
    ]
    wants_ranking = (
        any(p in query_lower for p in rank_phrases)
        or bool(re.search(r'\btop\s+\d+\b', query_lower))
        or bool(re.search(r'\btop\s+(five|ten|three)\b', query_lower))
    )
    if not wants_ranking:
        return None

    return 'crop' if crop_subject else 'product'


# Ask for as much as the index will return in one call.
#
# This was briefly lowered to 1000 on the theory that Pinecone caps top_k
# harder when metadata is included. That was wrong for this deployment and
# caused a real regression: a zero vector carries no similarity signal, so
# a smaller top_k returns an arbitrary subset rather than a representative
# one — and the 1000 that came back happened to contain no crop/product
# tags at all, turning "which crop generated the most complaints" into
# "no crop tags were found". Measured behaviour here is ~2300 records
# returned at top_k=10000, so keep it high.
#
# `complete` is still reported, because if the result set ever DOES hit the
# ceiling then every count derived from it is a lower bound and the reply
# must say so rather than presenting it as a census.
AGGREGATION_PAGE_SIZE = 10000


_WORD_NUMBERS = {
    "three": 3, "four": 4, "five": 5, "six": 6, "seven": 7,
    "eight": 8, "nine": 9, "ten": 10, "twenty": 20,
}


def detect_requested_top_n(query_lower: str) -> int | None:
    """Parse the N in "top 5" / "top five". detect_aggregation_request
    already recognises these phrasings, but the resolver hard-coded 10, so
    asking for the top 5 returned ten rows."""
    m = re.search(r'\btop\s+(\d{1,2})\b', query_lower)
    if m:
        n = int(m.group(1))
        return n if 1 <= n <= 50 else None
    m = re.search(r'\btop\s+(three|four|five|six|seven|eight|nine|ten|twenty)\b', query_lower)
    if m:
        return _WORD_NUMBERS[m.group(1)]
    return None


def fetch_matches_for_aggregation(index, filter_conditions, top_k=None):
    """Broad, non-semantic fetch used purely for exact counting over metadata
    tags. Returns (matches, complete) — `complete` is False when the fetch
    hit its ceiling, meaning any count derived from it is a lower bound."""
    matches, complete = _fetch_aggregation_page(index, filter_conditions, top_k)
    return matches, complete


def _fetch_aggregation_page(index, filter_conditions, top_k=None):
    page = top_k or AGGREGATION_PAGE_SIZE
    dummy_vector = [0.0] * EMBEDDING_DIMENSION
    # Deliberately NOT wrapped in try/except: a swallowed error here is
    # indistinguishable from "the dataset genuinely has nothing", which is
    # the worst possible way for a counting path to fail. Let it propagate
    # so the caller reports it with a reference id.
    results = index.query(
        vector=dummy_vector, top_k=page, include_metadata=True,
        filter=filter_conditions if filter_conditions else None
    )
    raw = results.get("matches", [])
    matches = [m for m in raw if not (m.get("metadata") or {}).get("is_stats_record")]
    # A full page means the result set was cut off, so any count derived
    # from it is a lower bound rather than a total.
    complete = len(raw) < page
    return matches, complete


def rank_by_field(matches, field: str, top_n: int = 10):
    """Count occurrences of each comma-separated tag value in the given metadata field, most common first."""
    counter = Counter()
    for m in matches:
        raw = m.get("metadata", {}).get(field, "")
        if not raw:
            continue
        for val in str(raw).split(","):
            val = val.strip()
            if val:
                counter[val] += 1
    return counter.most_common(top_n)


def detect_output_format(query_lower: str) -> str | None:
    """Detects which presentation format the user explicitly asked for."""
    if re.search(r'\bexcel\b|\bexport\b|\bdownload\b', query_lower):
        return 'excel'
    if re.search(r'\bppt\b|\bpowerpoint\b|\bpresentation\b|\bslides?\b', query_lower):
        return 'ppt'
    if re.search(r'\bexecutive\s+summary\b|\bone[- ]page\b|\bexec\s+summary\b', query_lower):
        return 'exec_summary'
    if re.search(r'\btable\b', query_lower):
        return 'table'
    if re.search(r'\bchart\b|\bgraph\b|\bvisuali[sz]e?\b|\bvisuali[sz]ation\b', query_lower):
        return 'chart'
    return None


def detect_trend_request(query_lower: str) -> bool:
    """ Detects requests for a monthly trend / growth-over-time breakdown (as opposed to a single-period or comparison-of-two-periods question). Deliberately conservative — bare words like "sales" or "products" alone are far too common in ordinary questions to trigger a full trend workup, so this requires an explicit temporal/trend framing. """
    trend_phrases = [
        'trend', 'trends', 'trending', 'month over month', 'month-over-month',
        'monthly trend', 'over time', 'growth', 'analytics over time',
        'performance over', 'monthly totals', 'monthly breakdown'
    ]
    return any(p in query_lower for p in trend_phrases)


def compute_monthly_trend(matches):
    """ Group raw Pinecone matches by (year, month) and return a chronologically sorted list of (label, count) tuples, e.g. [("January 2026", 42), ("February 2026", 51), ...]. Purely deterministic counting — no LLM involved, so the numbers are always exactly what's in the data. """
    counts = Counter()
    for m in matches:
        md = m.get("metadata", {})
        month = md.get("month")
        year = md.get("year")
        if not month or not year or month not in MONTH_ORDER:
            continue
        counts[(year, month)] += 1

    ordered_keys = sorted(counts.keys(), key=lambda k: (int(k[0]), MONTH_ORDER[k[1]]))
    return [(f"{month} {year}", counts[(year, month)]) for year, month in ordered_keys]


def densify_monthly_counts(monthly_counts):
    """Insert zero-count entries for months with no records.

    compute_monthly_trend only emits months that have data, so walking the
    result by index treats non-adjacent months as consecutive — a series of
    [Nov 2025, Feb 2026] produced a "+300% month-over-month" label across a
    three-month gap.
    """
    if not monthly_counts:
        return []
    parsed = []
    for label, count in monthly_counts:
        month, year = label.rsplit(" ", 1)
        parsed.append((_month_idx(month, year), count))
    out = []
    for idx in range(parsed[0][0], parsed[-1][0] + 1):
        month, year = _idx_to_month_year(idx)
        match = next((c for i, c in parsed if i == idx), 0)
        out.append((f"{month} {year}", match))
    return out


def compute_growth_series(monthly_counts):
    """ Given [(label, count), ...] in chronological order, return a parallel list of month-over-month growth percentages (None for the first period, which has no prior month to compare against). """
    growth = [None]
    for i in range(1, len(monthly_counts)):
        prev_count = monthly_counts[i - 1][1]
        curr_count = monthly_counts[i][1]
        if prev_count == 0:
            growth.append(None)
        else:
            growth.append(round((curr_count - prev_count) / prev_count * 100, 1))
    return growth


def build_pptx_report(title, subtitle, exec_summary_lines, kpis, chart_title, chart_labels, chart_values,
                       chart_type="column", table_headers=None, table_rows=None,
                       insights=None, recommendations=None):
    """ Build a professional PPTX deck: Title -> Executive Summary -> Key KPIs -> Chart -> Table -> Key Insights & Recommendations. Uses PowerPoint's native chart/table shapes (not an embedded image) so the deck stays fully editable in PowerPoint. Returns raw bytes. """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    # ── Slide 2: Executive Summary ──
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    lines = exec_summary_lines or ["No summary available."]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)

    # ── Slide 3: Key KPIs ──
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Key KPIs"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    kpi_items = list((kpis or {}).items()) or [("No KPIs available", "")]
    for i, (k, v) in enumerate(kpi_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.font.size = Pt(20)
        p.font.bold = True

    # ── Slide 4: Chart ──
    if chart_labels and chart_values:
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        title_box.text_frame.text = chart_title
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True

        chart_data = CategoryChartData()
        chart_data.categories = chart_labels
        chart_data.add_series(chart_title, chart_values)
        xl_type = XL_CHART_TYPE.LINE_MARKERS if chart_type == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
        slide.shapes.add_chart(xl_type, Inches(0.75), Inches(1.0), Inches(11.8), Inches(6.0), chart_data)

    # ── Slide 5: Table ──
    if table_headers and table_rows:
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        title_box.text_frame.text = "Supporting Data"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True

        MAX_ROWS = 15
        rows_to_show = table_rows[:MAX_ROWS]
        n_rows = len(rows_to_show) + 1
        n_cols = len(table_headers)
        table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0))
        table = table_shape.table
        for c, header in enumerate(table_headers):
            table.cell(0, c).text = str(header)
        for r, row in enumerate(rows_to_show, start=1):
            for c, val in enumerate(row):
                cell_text = str(val)
                table.cell(r, c).text = cell_text[:300]

    # ── Slide 6: Key Insights & Recommendations ──
    if insights or recommendations:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "Key Insights & Recommendations"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        first = True
        for point in (insights or []):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"Insight: {point}"
            p.font.size = Pt(16)
        for point in (recommendations or []):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"Recommendation: {point}"
            p.font.size = Pt(16)
            p.font.bold = True

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_excel_bytes(rows: list[dict]) -> bytes:
    """Turn a list of flat dict rows into an .xlsx file's raw bytes."""
    buf = BytesIO()
    pd.DataFrame(rows).to_excel(buf, index=False, engine="openpyxl")
    return buf.getvalue()


def build_csv_bytes(rows: list[dict], columns: list[str] | None = None) -> bytes:
    df = pd.DataFrame(rows, columns=columns) if columns else pd.DataFrame(rows)
    return df.to_csv(index=False).encode("utf-8")


def split_into_points(text: str, max_points: int = 6) -> list[str]:
    """Break an LLM prose/bullet response into short standalone points for slide bullets."""
    lines = [l.strip(" -*").strip() for l in text.split("\n") if l.strip(" -*").strip()]
    if len(lines) >= 2:
        return lines[:max_points]
    # Fall back to sentence-splitting for single-paragraph prose responses.
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    return [s.strip() for s in sentences if s.strip()][:max_points]


def build_subject_label(active_product, active_crop):
    """Combine crop + product into one display label, e.g. 'Wheat + Isabion'. Collapses to one when both detections landed on the same word (e.g. dynamic product-probe fallback re-matching the crop name)."""
    if active_crop and active_product and active_crop.lower() == active_product.lower():
        return active_crop.title()
    parts = []
    if active_crop:
        parts.append(active_crop.title())
    if active_product:
        parts.append(active_product.title())
    return " + ".join(parts) if parts else None


def build_header(query_intent, timeframe_label, active_product, periods, active_crop=None, category_filter=None):
    """ Product/crop and comparison context always take priority over the generic 'period' heading — a product or crop query is labeled with its subject (never falls back to a generic 'sentiment overview for the period' heading), and a comparison query is clearly labeled as a comparison. """
    subject_label = build_subject_label(active_product, active_crop)

    if periods:
        period_join = " 🆚 ".join(p[0] for p in periods)
        subject = f"{subject_label} — " if subject_label else ""
        if query_intent == "complaint":
            return f"🔀 {subject}Complaints Comparison: {period_join}\n\n"
        elif query_intent == "positive":
            return f"🔀 {subject}Positive Feedback Comparison: {period_join}\n\n"
        elif query_intent == "suggestion":
            return f"🔀 {subject}Suggestions Comparison: {period_join}\n\n"
        else:
            return f"🔀 {subject}Sentiment Comparison: {period_join}\n\n"

    if subject_label:
        suffix = f" ({timeframe_label})" if timeframe_label != DEFAULT_TIMEFRAME_LABEL else ""
        if query_intent == "complaint":
            return f"🐛 Complaints about {subject_label}{suffix}:\n\n"
        elif query_intent == "positive":
            return f"🌻 Positive Feedback about {subject_label}{suffix}:\n\n"
        elif query_intent == "suggestion":
            return f"💡 Suggestions about {subject_label}{suffix}:\n\n"
        elif query_intent == "topics":
            return f"🗣️ What's Being Said About {subject_label}{suffix}:\n\n"
        else:
            return f"🌾 {subject_label} — Sentiment Overview{suffix}:\n\n"

    if query_intent == "complaint":
        return f"🐛 Complaints of {timeframe_label}:\n\n"
    elif query_intent == "positive":
        return f"🌻 Positive Feedback of {timeframe_label}:\n\n"
    elif query_intent == "suggestion":
        return f"💡 Suggestions & Improvement Ideas for {timeframe_label}:\n\n"
    elif query_intent == "topics":
        return f"🗣️ Top Topics of {timeframe_label}:\n\n"
    elif category_filter == PRODUCT_QUERY_CATEGORY:
        return f"💰 Product Inquiries of {timeframe_label}:\n\n"
    else:
        return f"🌾 Sentiments of {timeframe_label}:\n\n"


def build_intent_badge(query_intent, active_product, periods, active_crop=None, category_filter=None):
    """Small colored pill label for the answer — purely cosmetic metadata, UI decides how to render it."""
    subject_label = build_subject_label(active_product, active_crop)
    if periods:
        return "🔀 Comparison"
    if subject_label:
        return f"🏷️ {subject_label}"
    if query_intent == "complaint":
        return "🐛 Complaints"
    if query_intent == "positive":
        return "🌻 Positive"
    if query_intent == "suggestion":
        return "💡 Suggestions"
    if query_intent == "topics":
        return "🗣️ Top Topics"
    if category_filter == PRODUCT_QUERY_CATEGORY:
        return "💰 Product Inquiries"
    return "🌾 Sentiment Overview"


def build_system_prompt(query_intent, timeframe_label, explicit_list_format, active_product, periods, active_crop=None, output_format=None, wants_products_only=False, avoid_repeat_text=None):
    """ Unified prompt builder. Preserves the original prose behaviour (including the two-paragraph favorable/complaints structure for the default sentiment case) while adding: real markdown bullet formatting when the user explicitly asks to "list" something, strict single-product/crop focus, explicit period-by-period comparison instructions, and output-format overrides (table / executive summary / PPT outline). """
    product_label = active_product.title() if active_product else None
    crop_label = active_crop.title() if active_crop else None

    if active_product and active_crop:
        product_clause = (
            f"Focus EXCLUSIVELY on the product '{product_label}' as it relates to the "
            f"crop '{crop_label}'. Do NOT mention any other product or crop, even if "
            f"they appear in the data context.\n"
        )
    elif active_product:
        product_clause = (
            f"Focus EXCLUSIVELY on the product '{product_label}'. Do NOT mention, "
            f"reference, or summarize information about any other product, even if "
            f"other products appear in the data context — ignore anything not about "
            f"'{product_label}'.\n"
        )
    elif active_crop:
        product_clause = (
            f"Focus EXCLUSIVELY on the crop '{crop_label}'. Do NOT mention, reference, "
            f"or summarize information about any other crop, even if other crops appear "
            f"in the data context — ignore anything not about '{crop_label}'. Still name "
            f"every product mentioned in connection with '{crop_label}'.\n"
        )
    else:
        product_clause = (
            "Explicitly name every product mentioned in the data context along with "
            "the exact reason for the feedback.\n"
        )

    comparison_clause = ""
    if periods:
        period_names = ", ".join(p[0] for p in periods)
        comparison_clause = (
            f"This is a COMPARISON request across these periods: {period_names}. "
            f"The data context below is divided into clearly labeled sections, one per "
            f"period. Explicitly compare the periods against each other — call out what "
            f"increased, decreased, improved, worsened, or stayed roughly the same. "
            f"Refer to each period by its exact name. "
            f"CRITICAL: for every point you make, name the specific product it is about "
            f"(never speak only in generic sentences with no product named), and for each "
            f"period state plainly whether that product's feedback was positive/satisfactory "
            f"or negative/unsatisfactory in that period — e.g. 'In {period_names.split(', ')[0]}, "
            f"growers were satisfied with <Product>, but in the other period they were not.' "
            f"Do this for every product that appears in the data context.\n"
        )

    intent_label = {
        "complaint":  "complaints and concerns (including root-cause issues)",
        "positive":   "positive feedback and appreciation",
        "suggestion": "grower suggestions and improvement recommendations",
        "sentiment":  "overall sentiment (both positive and negative)",
        "topics":     "the most frequently discussed topics and themes",
    }[query_intent]

    opening_hint = {
        "complaint":  f"e.g. 'Here are the complaints for {timeframe_label}:'",
        "positive":   f"e.g. 'The positive feedback for {timeframe_label} looks great!'",
        "suggestion": f"e.g. 'Here are the grower suggestions for {timeframe_label}:'",
        "sentiment":  f"e.g. 'Here is the sentiment overview for {timeframe_label}:'",
        "topics":     f"e.g. 'Here's what growers are talking about most for {timeframe_label}:'",
    }[query_intent]
    subject_label = build_subject_label(active_product, active_crop)
    if subject_label:
        opening_hint = f"e.g. 'Here's what growers are saying about {subject_label}:'"

    output_format_clause = ""
    if output_format == "table":
        output_format_clause = (
            "OUTPUT FORMAT OVERRIDE: Format your ENTIRE response as a markdown table "
            "with columns '| Category | Feedback |'. One data point per row. No prose "
            "outside the table.\n"
        )
    elif output_format == "exec_summary":
        output_format_clause = (
            "OUTPUT FORMAT OVERRIDE: Write a one-page executive summary with these "
            "bold section headers on their own lines, in order: '*Headline:*' (one "
            "sentence takeaway), '*Key Insights:*' (3-5 bullet points, one specific "
            "point each), '*Recommendation:*' (1-2 sentences on what to do next).\n"
        )
    elif output_format == "ppt":
        output_format_clause = (
            "OUTPUT FORMAT OVERRIDE: Format your response as a PowerPoint-ready slide "
            "outline. Use '*Slide 1: <short title>*' on its own line, followed by "
            "3-5 short bullet points (each starting with '- '), then a blank line "
            "before the next slide if more than one slide is warranted. Keep every "
            "bullet short enough to fit on a slide (under 15 words).\n"
        )

    if explicit_list_format:
        format_clause = (
            "Format your ENTIRE response as a real markdown bullet list. Every single "
            "bullet MUST start on its own new line with a dash and a space: '- '. "
            "Never write '•' and never put more than one bullet on the same line. "
            "Each bullet must be one specific, concrete point (one product/issue per "
            "bullet) — no paragraphs, no prose outside the list.\n"
        )
        if query_intent == "sentiment" and not periods:
            format_clause += (
                "Group the bullets under two bold headers on their own lines: "
                "'*Positive:*' followed by positive bullets, then a blank line, then "
                "'*Negative:*' followed by negative bullets.\n"
            )
        if periods:
            format_clause += (
                "Group the bullets under one bold header per period (using the exact "
                "period names given above, each on its own line), followed by that "
                "period's bullets, with a blank line between groups.\n"
            )
        structure_clause = ""
    else:
        format_clause = (
            "Respond in natural, flowing prose — no bullet points, no markdown lists, "
            "no asterisks. Sound like a helpful chatbot, not a formal report. Keep it "
            "concise.\n"
        )
        if query_intent == "sentiment" and not periods and not active_product:
            structure_clause = (
                "Structure your response in exactly two short paragraphs:\n"
                "Paragraph 1 — Favorable Sentiments: summarize positive trends.\n"
                "Paragraph 2 — Complaints & Concerns: summarize issues.\n"
                "Each paragraph should be 3-5 sentences max.\n"
            )
        elif query_intent == "topics":
            structure_clause = (
                "Structure your response as one short paragraph per topic (3-6 "
                "topics total, most-discussed first), each naming the topic and "
                "summarizing what's said about it in 1-2 sentences.\n"
            )
        else:
            structure_clause = "Keep the response to 4-6 sentences max.\n"

    topics_clause = ""
    if query_intent == "topics":
        topics_clause = (
            "This is a TOPICS/THEMES request, not a sentiment request: identify "
            "the 3-6 most frequently occurring topics or themes in the data "
            "context below (e.g. pricing, product performance, packaging, "
            "availability, application timing) and briefly describe what "
            "growers are saying about each one. Do NOT frame this as "
            "positive-vs-negative sentiment analysis — focus on WHAT is being "
            "discussed, not whether it is good or bad. Order topics from most "
            "to least frequently mentioned where you can tell.\n"
        )

    # An explicit output-format request (table / exec summary / PPT outline)
    # always wins over the default bullet/prose formatting instructions.
    if output_format_clause:
        format_clause = output_format_clause
        structure_clause = ""

    system_prompt = (
        "You are a smart, friendly chatbot analyst for Syngenta, an agriculture company. "
        "STRICT GROUNDING RULE — READ CAREFULLY: You must use ONLY the information given "
        "to you in the 'Data Context' block in the user's message. You have general "
        "knowledge about real Syngenta/agriculture products from your training — you must "
        "IGNORE all of that here. Do NOT invent, assume, guess, or add any product name, "
        "complaint, statistic, or feedback point that is not explicitly written in the Data "
        "Context, even if it sounds plausible or matches a real product you know about. If "
        "the Data Context contains only one point, your entire response must be based on "
        "that single point only — never pad the list with extra products or details to make "
        "it look longer or more complete. If the Data Context is empty or has nothing "
        "relevant, say so plainly instead of making something up. "
        f"Cover ONLY {intent_label} from the data context provided. "
        f"{product_clause}"
        + (
            "The user is asking specifically about PRODUCTS. Only name real "
            "Syngenta product/brand names (e.g. Isabion, Axial, Quantis, "
            "Cropwise). Do NOT list crop diseases, pests, weeds, or agronomic "
            "problems (e.g. Septoria, Blast, Armyworm, termites, weeds, "
            "yellowing) as if they were products — mention a disease/pest "
            "only in passing if it explains why a product was used, never as "
            "an item in a products list.\n"
            if wants_products_only else ""
        )
        + f"{comparison_clause}"
        f"{topics_clause}"
        f"{format_clause}"
        f"{structure_clause}"
        f"Start your response with a short, clear opening line ({opening_hint}), then "
        "continue. Write so a busy reader understands the key takeaway at first glance. "
        "Do not include bracketed dates, week labels, or raw metadata tags in the output. "
        "REMINDER: every product name and every point in your response must come directly "
        "from the Data Context above — never introduce a product or detail that isn't "
        "explicitly there."
    )

    if avoid_repeat_text:
        system_prompt += (
            "\n\nCONTINUATION REQUEST: The user already received the response below "
            "earlier in this conversation and is now explicitly asking for MORE or "
            "DIFFERENT points on the same subject — not a repeat. Do not restate any "
            "point from it in substantially the same words. Pull NEW points from the "
            "Data Context that weren't already covered. If nothing further is "
            "supported by the Data Context, say plainly that there isn't anything "
            "further to add rather than repeating the earlier points.\n\n"
            f"Previous response already given:\n{avoid_repeat_text}"
        )

    return system_prompt


# ==========================================
# EXCEL INGESTION (no UI calls — raises on error, returns a summary dict)
# ==========================================

def _make_metadata_payload(inferred_year, row_month, week_label, category, bullet,
                           sheet_name="", src_row=-1, ingest_run="", week_num=None):
    """Build the embedding text and the Pinecone metadata for one bullet.

    Provenance (sheet / src_row / ingest_run) is stored so a re-ingest can
    delete exactly what it is replacing. Without it there was no way to
    scope a delete, so corrections could never take effect: the old vector
    simply stayed in the index forever.

    The old "text" field duplicated the whole bullet inside a longer
    sentence and was stored alongside "value", roughly doubling metadata
    size for no benefit — a single long cell could push a payload past
    Pinecone's 40KB limit and reject the entire 50-vector batch. It is
    reconstructible from the other fields, so it is no longer stored.
    """
    is_positive = category in POSITIVE_CATEGORIES
    is_negative = category in NEGATIVE_CATEGORIES
    context_chunk = (
        f"Year: {inferred_year}. "
        f"Month: {row_month}. "
        f"Week: {week_label}. "
        f"Case Category: {category}. "
        f"Feedback: {bullet}."
    )
    metadata = {
        "month":     row_month,
        "year":      inferred_year,
        "week":      week_label,
        "category":  category,
        "sentiment": (
            "positive" if is_positive
            else "negative" if is_negative
            else "neutral"
        ),
        "value":    bullet[:MAX_VALUE_CHARS],
        "crop":     ",".join(extract_crops(bullet)),
        "products": ",".join(extract_product_mentions(bullet)),
        "sheet":    sheet_name,
        "src_row":  int(src_row),
        "ingest_run": ingest_run,
    }
    if week_num is not None:
        # Stored as an integer so week can be a database-level filter
        # instead of a Python substring test applied after the top_k cut.
        metadata["week_num"] = int(week_num)
    if len(bullet) > MAX_VALUE_CHARS:
        metadata["truncated"] = True
    return context_chunk, metadata


def _vector_id(sheet: str, category: str, week_label: str, row: int, bullet_idx: int, bullet: str) -> str:
    """Content-addressed id.

    The previous scheme concatenated row and bullet indexes with no
    separator ("...{idx}{b_idx}"), so row 1/bullet 12 and row 11/bullet 2
    both produced "...112" and the second silently overwrote the first —
    16 records reported ingested, 13 actually stored. Hashing the content
    also makes re-ingesting unchanged rows idempotent.
    """
    basis = f"{sheet}|{category}|{week_label}|{row}|{bullet_idx}|{bullet}"
    return "v_" + hashlib.sha1(basis.encode("utf-8")).hexdigest()


def _clean_cell_text(cell_val) -> str | None:
    """Return usable feedback text, or None if this cell isn't feedback.

    Guards against non-text cells: a date cell stringified to
    '2026-01-05 00:00:00' and a float to '12.0', both of which were being
    embedded and counted as grower feedback.
    """
    if cell_val is None:
        return None
    if isinstance(cell_val, (int, float, bool)):
        return None
    text = str(cell_val).strip()
    if not text or is_empty_cell(text):
        return None
    if re.fullmatch(r'[\d.\-:/ ]+', text):
        return None
    if re.fullmatch(r'\d{4}-\d{2}-\d{2}[ T].*', text):
        return None
    return text


def _week_number(week_label: str) -> int | None:
    """Parse the ordinal from a week label ('2nd Week January' -> 2).

    Matches the ordinal specifically rather than the first integer in the
    string — a column named 'March 2026 Week 1' previously yielded 2026.
    """
    m = re.search(r'\b(\d{1,2})\s*(?:st|nd|rd|th)\b', week_label, re.IGNORECASE)
    if m:
        n = int(m.group(1))
        return n if 1 <= n <= 6 else None
    m = re.search(r'\bweek\s*(\d{1,2})\b', week_label, re.IGNORECASE)
    if m:
        n = int(m.group(1))
        return n if 1 <= n <= 6 else None
    return None


_WEEK_COL_RE = re.compile(r'^w(?:ee)?k(?:\s*(?:no\.?|number|#))?$', re.IGNORECASE)


def run_ingestion(file_bytes: bytes, pinecone_api_key: str, purge_first: bool = False) -> dict:
    """Parse the workbook, tag each bullet, embed, and upsert into Pinecone.

    Returns {"total_records", "summary", "skipped", "ingest_run"}. The
    "skipped" list is the important addition: previously four separate
    silent `continue` paths could drop entire sheets while the call still
    reported success, so a workbook could half-ingest with no indication.

    purge_first=True deletes the whole index before writing. Required once
    after any change to the tagging logic, because vectors already in the
    index keep their old (wrong) metadata — re-ingesting alone does not
    repair them.
    """
    excel_file = pd.ExcelFile(BytesIO(file_bytes))
    all_sheets = excel_file.sheet_names

    pc = Pinecone(api_key=pinecone_api_key)
    index = pc.Index(PINECONE_INDEX_NAME)

    ingest_run = hashlib.sha1(f"{len(file_bytes)}|{','.join(map(str, all_sheets))}".encode()).hexdigest()[:12]

    payload_batch = []
    text_inputs_for_embedding = []
    discovered_data_summary = {}
    skipped: list[dict] = []

    def skip(sheet, reason, detail=""):
        skipped.append({"sheet": str(sheet), "reason": reason, "detail": str(detail)})

    for sheet_name in all_sheets:
        sheet_clean = str(sheet_name).strip()
        inferred_year = infer_year_for_sheet(sheet_clean, [str(s) for s in all_sheets])
        if not inferred_year:
            skip(sheet_clean, "no_year", "Could not infer a year for this sheet, or the inferred year would duplicate an explicitly dated sheet.")
            continue

        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        df.columns = [re.sub(r'\s+', ' ', str(c)).strip() for c in df.columns]

        # Duplicate labels make row[col] return a Series, whose repr
        # ('Name: 0, dtype: str') was being ingested as feedback text.
        if pd.Index(df.columns).duplicated().any():
            dupes = sorted({c for c in df.columns if list(df.columns).count(c) > 1})
            skip(sheet_clean, "duplicate_columns", f"Columns collide after whitespace normalization: {dupes}")
            continue

        cat_col = find_category_column(df.columns)

        if cat_col:
            # ── LAYOUT A: categories are ROWS, weeks are COLUMNS ──
            week_cols = [c for c in df.columns if 'week' in c.lower()]
            if not week_cols:
                skip(sheet_clean, "no_week_columns", "Layout A detected (a category column exists) but no column header contains 'week'.")
                continue

            # Merged category cells give the value only to the first row;
            # every continuation row came back NaN and was dropped. Layout B
            # already forward-filled month/week; Layout A filled nothing.
            df[cat_col] = df[cat_col].replace(r'^\s*$', pd.NA, regex=True).ffill()

            unmapped = set()
            for idx, row in df.iterrows():
                category = normalize_category(row.get(cat_col, None))
                if not category:
                    raw = str(row.get(cat_col, "")).strip()
                    if raw and raw.lower() not in EMPTY_VALUES:
                        unmapped.add(raw)
                    continue

                for col in week_cols:
                    cell_raw = _clean_cell_text(row[col])
                    if cell_raw is None:
                        continue

                    bullets = split_bullets(cell_raw)
                    if not bullets:
                        continue

                    row_month = extract_month_from_col(col)
                    if row_month == "Unknown":
                        skip(sheet_clean, "unparseable_month", f"Week column '{col}' contains no month name; its rows are unreachable by any month filter and were skipped.")
                        continue

                    stat_key = f"{row_month} {inferred_year}"
                    discovered_data_summary[stat_key] = (
                        discovered_data_summary.get(stat_key, 0) + len(bullets)
                    )

                    for b_idx, bullet in enumerate(bullets):
                        context_chunk, metadata_payload = _make_metadata_payload(
                            inferred_year, row_month, col, category, bullet,
                            sheet_name=sheet_clean, src_row=idx, ingest_run=ingest_run,
                            week_num=_week_number(col),
                        )
                        payload_batch.append({
                            "id": _vector_id(sheet_clean, category, col, idx, b_idx, bullet),
                            "metadata": metadata_payload,
                        })
                        text_inputs_for_embedding.append(context_chunk)

            if unmapped:
                skip(sheet_clean, "unmapped_categories", f"These category values are not recognized and their rows were skipped: {sorted(unmapped)}")

        else:
            # ── LAYOUT B: Month/Week are ROW values, categories are COLUMN headers ──
            df_raw = pd.read_excel(excel_file, sheet_name=sheet_name, header=None)

            header_row_idx = None
            for i in range(min(10, len(df_raw))):
                row_vals = [re.sub(r'\s+', ' ', str(v)).strip().lower() for v in df_raw.iloc[i].tolist()]
                if any(v == 'month' for v in row_vals):
                    header_row_idx = i
                    break
            if header_row_idx is None:
                skip(sheet_clean, "no_header_row", "Layout B expected but no row in the first 10 contains a 'Month' header cell.")
                continue

            # Forward-fill the header row so a merged header spanning two
            # columns doesn't silently discard the second column's data.
            header_cells = df_raw.iloc[header_row_idx].tolist()
            col_map, last_label = {}, None
            for j, v in enumerate(header_cells):
                text = re.sub(r'\s+', ' ', str(v)).strip()
                if text and text.lower() != 'nan':
                    last_label = text
                    col_map[j] = text
                elif last_label is not None:
                    col_map[j] = last_label

            month_col_idx = next((i for i, v in col_map.items() if v.strip().lower() == 'month'), None)
            week_col_idx = next((i for i, v in col_map.items() if _WEEK_COL_RE.match(v.strip())), None)

            # Only columns whose header resolves to a KNOWN category are
            # feedback. Previously every non-Month/Week column qualified, so
            # Region / Remarks / Dealer values were ingested as feedback and
            # region names were counted as products.
            category_cols = {}
            rejected_headers = []
            for i, v in col_map.items():
                if i in (month_col_idx, week_col_idx):
                    continue
                if normalize_category(v):
                    category_cols[i] = v
                else:
                    rejected_headers.append(v)

            if month_col_idx is None or not category_cols:
                skip(sheet_clean, "no_category_columns", f"No column header maps to a known feedback category. Headers seen: {sorted(set(col_map.values()))}")
                continue
            if rejected_headers:
                skip(sheet_clean, "ignored_columns", f"Not feedback categories, so not ingested: {sorted(set(rejected_headers))}")

            current_month = None
            current_week = None
            for r in range(header_row_idx + 1, len(df_raw)):
                row = df_raw.iloc[r]

                mval = row[month_col_idx]
                if pd.notna(mval) and str(mval).strip():
                    parsed = extract_month_from_col(str(mval).strip())
                    if parsed == "Unknown":
                        # A 'Total'/'Notes' row must not overwrite the
                        # forward-filled month for every row beneath it.
                        continue
                    if parsed != current_month:
                        # Reset the week when the month changes, or a blank
                        # week cell on a month's first row inherits the
                        # previous month's last week.
                        current_week = None
                    current_month = parsed

                if week_col_idx is not None:
                    wval = row[week_col_idx]
                    if pd.notna(wval) and str(wval).strip():
                        current_week = str(wval).strip()

                if not current_month:
                    continue

                week_label = f"{current_week} Week {current_month}" if current_week else current_month

                for col_idx, raw_category_name in category_cols.items():
                    category = normalize_category(raw_category_name)
                    if not category:
                        continue

                    cell_raw = _clean_cell_text(row[col_idx])
                    if cell_raw is None:
                        continue

                    bullets = split_bullets(cell_raw)
                    if not bullets:
                        continue

                    stat_key = f"{current_month} {inferred_year}"
                    discovered_data_summary[stat_key] = (
                        discovered_data_summary.get(stat_key, 0) + len(bullets)
                    )

                    for b_idx, bullet in enumerate(bullets):
                        context_chunk, metadata_payload = _make_metadata_payload(
                            inferred_year, current_month, week_label, category, bullet,
                            sheet_name=sheet_clean, src_row=r, ingest_run=ingest_run,
                            week_num=_week_number(week_label),
                        )
                        payload_batch.append({
                            "id": _vector_id(sheet_clean, category, week_label, r, b_idx, bullet),
                            "metadata": metadata_payload,
                        })
                        text_inputs_for_embedding.append(context_chunk)

    # Collapse any duplicate ids before embedding, so the reported count
    # matches what actually lands in the index.
    deduped, seen_ids = [], set()
    deduped_texts = []
    for item, text in zip(payload_batch, text_inputs_for_embedding):
        if item["id"] in seen_ids:
            continue
        seen_ids.add(item["id"])
        deduped.append(item)
        deduped_texts.append(text)
    payload_batch, text_inputs_for_embedding = deduped, deduped_texts

    total_records = len(payload_batch)
    if total_records == 0:
        raise ValueError(
            "No records found in the uploaded workbook."
            + (f" Skipped: {skipped}" if skipped else "")
        )

    if purge_first:
        # Vectors already in the index keep their old metadata, and the
        # content-hash ids of corrected rows differ from the positional ids
        # they replace — so without a purge the stale records survive.
        try:
            index.delete(delete_all=True)
        except Exception as e:
            raise RuntimeError(f"Purge requested but the index delete failed, so ingestion was aborted to avoid mixing old and new data: {e}")

    # Embed and upsert in the SAME pass. Collecting every embedding first
    # meant peak memory held the whole dataset (~12KB/record), which OOMs a
    # 512MB instance well before a realistic workbook is finished.
    BATCH_LIMIT = 96
    written = 0
    for i in range(0, total_records, BATCH_LIMIT):
        text_batch = text_inputs_for_embedding[i: i + BATCH_LIMIT]
        meta_batch = payload_batch[i: i + BATCH_LIMIT]

        embeddings_response = pc.inference.embed(
            model="llama-text-embed-v2",
            inputs=text_batch,
            parameters={"input_type": "passage", "dimension": EMBEDDING_DIMENSION}
        )
        values = [item.values for item in embeddings_response]
        if len(values) != len(meta_batch):
            # Positional pairing means a short response would silently
            # attach every later vector to the wrong metadata.
            raise RuntimeError(
                f"Embedding count mismatch ({len(values)} vectors for {len(meta_batch)} records); "
                f"aborted after {written} records to avoid writing mismatched data."
            )

        vectors = [
            {"id": m["id"], "values": v, "metadata": m["metadata"]}
            for m, v in zip(meta_batch, values)
        ]
        for j in range(0, len(vectors), 50):
            index.upsert(vectors=vectors[j: j + 50])
        written += len(vectors)

    # Record the dataset's real extents so relative-date resolution reads a
    # fact instead of sampling arbitrary records and taking a max.
    dated = [
        (int(m["metadata"]["year"]), m["metadata"]["month"])
        for m in payload_batch
        if str(m["metadata"].get("year", "")).isdigit() and m["metadata"].get("month") in MONTH_ORDER
    ]
    if dated:
        max_year, max_month = max(dated, key=lambda p: (p[0], MONTH_ORDER[p[1]]))
        min_year, min_month = min(dated, key=lambda p: (p[0], MONTH_ORDER[p[1]]))
        try:
            index.upsert(vectors=[{
                "id": INDEX_STATS_ID,
                "values": [0.0] * EMBEDDING_DIMENSION,
                "metadata": {
                    "is_stats_record": True,
                    "max_year": str(max_year), "max_month": max_month,
                    "min_year": str(min_year), "min_month": min_month,
                    "total_records": written,
                    "ingest_run": ingest_run,
                },
            }])
        except Exception:
            pass  # stats are an optimization; the sampled fallback still works

    return {
        "total_records": written,
        "summary": discovered_data_summary,
        "skipped": skipped,
        "ingest_run": ingest_run,
    }

# ==========================================
# CHAT ORCHESTRATION (two-phase: process → stream LLM yourself → finalize)
# ==========================================

def is_query_in_scope(user_query: str) -> bool:
    """Strict topic guardrail: at least one recognized domain word must appear."""
    query_words = re.findall(r'\b\w+\b', user_query.lower())
    return any(word in ALLOWED_GUARDRAIL_KEYWORDS for word in query_words)


FOLLOWUP_PHRASES = [
    "what about", "how about", "what's about", "and what about",
    "same for", "also for", "and for", "what of", "and about",
    "what if", "and how about",
]

# Phrases that explicitly ask for MORE/DIFFERENT points on the same subject
# rather than a repeat of the last answer — "what other insights", "anything
# else", "what more can you tell me". These also count as a follow-up
# continuation (so product/crop/intent are inherited), but additionally
# signal that the previous answer's exact text should be passed to the LLM
# as "already said" content to avoid parroting the same points back.
MORE_INSIGHTS_PHRASES = [
    "what other", "any other", "anything else", "what else",
    "else can you", "more insight", "other insight", "something else",
    "what more", "anything more", "else you can", "more detail",
    "more information", "additional insight", "additional detail",
]

FOLLOWUP_PHRASES = FOLLOWUP_PHRASES + MORE_INSIGHTS_PHRASES


def detect_followup_reference(query_lower: str) -> bool:
    """ Detects explicit conversational continuation phrasing ("what about wheat?", "and for Isabion?") — deliberately narrow (exact phrase match, not just "short query") so an unrelated fresh question never accidentally inherits stale context from a few turns ago. """
    return any(p in query_lower for p in FOLLOWUP_PHRASES)


def detect_wants_more(query_lower: str) -> bool:
    """Detects a request for MORE/DIFFERENT points on the same subject, as opposed to a plain repeat of the previous question — used to tell the LLM what was already said so it doesn't repeat itself."""
    return any(p in query_lower for p in MORE_INSIGHTS_PHRASES)


# Phrases that signal the user is correcting or giving feedback ABOUT the
# chatbot itself ("Kaho is not a product", "you're wrong about that") —
# rather than asking a question about the grower-feedback data. Deliberately
# a deterministic phrase list, not an LLM call: this class of message is
# rare enough that adding LLM latency/cost to every single message just to
# catch it isn't worth it, mirroring the cheap-keyword-list-first approach
# used for intent detection. Real natural language can phrase a correction
# in effectively unlimited ways, so — like the relative-date parser — this
# is a real, documented coverage limit, not a claim of completeness.
CORRECTION_PHRASES = [
    "is not a product", "isn't a product", "not a real product",
    "is not a real product", "is not a crop", "isn't a crop",
    "not a real crop", "that's wrong", "that is wrong", "that's incorrect",
    "that is incorrect", "that's not correct", "that is not correct",
    "you're wrong", "you are wrong", "you're mistaken", "you are mistaken",
    "you made a mistake", "you made an error", "this is incorrect",
    "this is wrong", "please fix this", "please correct this",
    "stop treating", "stop showing", "stop labeling", "stop calling",
    "should not be treated as", "shouldn't be treated as",
    "you got that wrong", "you got this wrong",
]

CORRECTION_ACK_REPLY = (
    "Thanks for flagging that — I don't have a way to update my product/crop "
    "catalog or behavior directly from a chat message, since that's shared "
    "across everyone using this tool and changing it here could affect "
    "results for other users. I've logged this note for the team to review "
    "and correct in the underlying configuration if confirmed."
)


def detect_correction_or_meta_feedback(query_lower: str) -> bool:
    """Detects a message that's giving feedback/correcting the chatbot's own behavior or knowledge, rather than asking a data question — see CORRECTION_PHRASES for the coverage caveat."""
    return any(p in query_lower for p in CORRECTION_PHRASES)


def process_chat_query(user_query: str, pinecone_api_key: str, groq_api_key: str | None = None, prior_context: dict | None = None) -> dict:
    """ Runs everything up to (but not including) the Groq call. Returns a dict with "kind": - "blocked": off-topic query. "reply" is ready to show. - "no_key": Pinecone not configured. "reply" is ready to show. - "ranking" / "trend" / "no_data": fully resolved without needing an LLM — "reply" (markdown), optionally "chart" ({"type","title","labels","values"}) and "downloads" ({"csv","excel","pptx"} bytes). - "normal": needs an LLM call. Contains "system_prompt" / "user_prompt" ready to send to Groq, plus everything finalize_normal_response() will need afterwards. prior_context (optional): {"product","crop","intent"} resolved from the previous turn. Only used to fill in slots the CURRENT query left unspecified, and only when the query contains an explicit continuation phrase ("what about wheat?") — never silently overrides anything the current query itself states, so a fresh unrelated question is never scoped by accident. """
    query_lower = user_query.lower()

    # Checked FIRST, before the topic guardrail: a correction like "you're
    # wrong about that" carries no domain vocabulary and would otherwise hit
    # the generic "I cannot generate this response" guardrail message,
    # which is just as unhelpful as running it through the data pipeline —
    # neither engages with what the user actually said.
    if detect_correction_or_meta_feedback(query_lower):
        return {"kind": "meta_feedback", "reply": CORRECTION_ACK_REPLY}

    # A bare continuation phrase ("what about last month?") carries no
    # domain keyword of its own — it only makes sense in light of a prior
    # in-scope turn, so it's allowed past the guardrail exactly when there
    # IS prior context for it to continue.
    is_followup_continuation = bool(prior_context) and detect_followup_reference(query_lower)

    if not is_query_in_scope(user_query) and not is_followup_continuation:
        return {
            "kind": "blocked",
            "reply": (
                "I cannot generate this response. "
                "I am strictly locked to analyzed dataset metrics "
                "and cannot find relevant information for this query."
            ),
        }

    if not pinecone_api_key:
        return {"kind": "no_key", "reply": "🤖 Execution Halted: Pinecone API key is not configured."}

    all_months = extract_all_months(query_lower)
    all_years = extract_all_years(query_lower)
    all_weeks = extract_all_weeks(query_lower)

    detected_month = all_months[0] if all_months else None
    detected_year = all_years[0] if all_years else None
    detected_week = all_weeks[0] if all_weeks else None

    wants_last_week = bool(re.search(r'\b(last|latest|most recent|recent)\s+week\b', query_lower))
    explicit_list_format = bool(re.search(r'\blist(ed|ing)?\b|\bbullets?\b|\bbullet\s*points?\b', query_lower))
    output_format = detect_output_format(query_lower)
    wants_products_only = bool(re.search(r'\bproducts?\b', query_lower))
    wants_trend = detect_trend_request(query_lower)
    aggregation_dimension = detect_aggregation_request(query_lower)
    requested_top_n = detect_requested_top_n(query_lower)

    topic_keywords = [
        "talking about", "talking most about", "talking the most about",
        "most discussed", "most talked about", "common themes",
        "main themes", "what topics", "which topics", "trending topics",
        "top topics", "hot topics", "being discussed", "being talked about",
        "conversation topics", "main subjects", "most common subjects",
    ]
    complaint_keywords = [
        "complaint", "complaints", "negative feedback",
        "negative", "issues", "problems", "concerns",
        "issue", "problem", "root cause", "root causes",
        # expanded — real phrasings that don't literally say "complaint"/"issue"
        "dissatisfied", "unhappy", "frustration", "frustrated",
        "unavailability", "shortage", "shortages", "delay", "delays",
        "delayed", "defect", "defects", "faulty", "damaged", "not working",
        "poor quality", "low quality", "bad experience", "disappointed",
        "difficulty", "difficulties", "trouble", "troubles",
        "pain point", "pain points", "grievance", "grievances",
    ]
    positive_keywords = [
        "positive feedback", "appreciation", "praise",
        "favorable", "satisfied",
        # expanded
        "happy", "pleased", "impressed", "love", "loved", "loving",
        "great experience", "good experience", "worked well", "works well",
        "effective", "satisfaction", "delighted", "thrilled",
    ]
    suggestion_keywords = [
        "suggestion", "suggestions", "recommend", "recommendation",
        "recommendations", "improvement", "improvements",
        "expectation", "expectations",
        # expanded
        "would like", "wish", "wishes", "hope for", "hoping for",
        "request", "requests", "requested", "ask for", "asking for",
        "want to see", "should add", "should include", "feature request",
        "enhancement", "enhancements", "could improve",
    ]
    sentiment_keywords = [
        "sentiment", "sentiments", "overall", "general",
        "overview", "analysis", "summary", "both",
        "feedback", "feedbacks",
        # expanded — general "what do people think" style asks
        "think", "thoughts", "opinion", "opinions", "views",
        "perception", "perceptions", "reaction", "reactions",
        "impression", "impressions", "experience", "experiences",
    ]

    query_intent = "sentiment"
    intent_explicit = False
    if any(phrase in query_lower for phrase in topic_keywords):
        # Checked first: "what are growers talking about" is a THEMES/TOPICS
        # question, not a positive/negative valence question — forcing it
        # through the same sentiment-overview template is exactly the
        # one-size-fits-all problem this intent exists to avoid.
        query_intent = "topics"
        intent_explicit = True
    elif any(phrase in query_lower for phrase in complaint_keywords):
        query_intent = "complaint"
        intent_explicit = True
    elif any(phrase in query_lower for phrase in positive_keywords):
        query_intent = "positive"
        intent_explicit = True
    elif any(phrase in query_lower for phrase in suggestion_keywords):
        query_intent = "suggestion"
        intent_explicit = True
    elif any(word in query_lower for word in sentiment_keywords):
        query_intent = "sentiment"
        intent_explicit = True

    category_filter = SUGGESTION_CATEGORY if query_intent == "suggestion" else None
    sales_scoped = False
    if not category_filter and query_intent == "sentiment" and any(k in query_lower for k in SALES_KEYWORDS):
        category_filter = PRODUCT_QUERY_CATEGORY
        sales_scoped = True

    pc = Pinecone(api_key=pinecone_api_key)
    index = pc.Index(PINECONE_INDEX_NAME)

    if not all_years:
        if re.search(r'\byear[- ]over[- ]year\b|\byoy\b', query_lower):
            latest = get_latest_year_from_index(index)
            try:
                all_years = [str(int(latest) - 1), latest]
            except ValueError:
                pass
        elif re.search(r'\blast\s+year\b', query_lower):
            latest = get_latest_year_from_index(index)
            try:
                all_years = [str(int(latest) - 1)]
            except ValueError:
                pass
        elif re.search(r'\bthis\s+year\b|\bcurrent\s+year\b', query_lower):
            all_years = [get_latest_year_from_index(index)]
        detected_year = all_years[0] if all_years else None

    try:
        query_response = pc.inference.embed(
            model="llama-text-embed-v2",
            inputs=[user_query],
            parameters={"input_type": "query", "dimension": EMBEDDING_DIMENSION}
        )
        query_vector = query_response[0].values
    except Exception as e:
        return {"kind": "no_data", "reply": f"Query embedding failed: {e}"}

    active_product = detect_product_known(query_lower)
    if not active_product:
        active_product = detect_product_dynamic(query_lower, index, pc, original_query=user_query)
    active_crop = detect_crop(query_lower)
    if active_product and active_crop and active_product.lower() == active_crop.lower():
        active_product = None

    # ── Follow-up slot inheritance: only fills gaps the CURRENT query left
    # unspecified, and only on an explicit continuation phrase — never
    # overrides anything the current query itself states. ──
    avoid_repeat_text = None
    if prior_context and detect_followup_reference(query_lower):
        if not active_product and prior_context.get("product"):
            active_product = prior_context["product"]
        if not active_crop and prior_context.get("crop"):
            active_crop = prior_context["crop"]
        if not intent_explicit and prior_context.get("intent"):
            query_intent = prior_context["intent"]
            category_filter = SUGGESTION_CATEGORY if query_intent == "suggestion" else category_filter
        # "what other insights...", "anything else..." — explicitly asking
        # for MORE/DIFFERENT points, not a repeat. Pass the previous answer
        # to the LLM so it doesn't parrot the same points back.
        if detect_wants_more(query_lower):
            avoid_repeat_text = prior_context.get("last_reply")

    # ── LLM-assisted fallback for genuinely ambiguous queries — only fires
    # when regex-based detection found NOTHING at all (no product, no crop,
    # no explicit intent, AND no sales/pricing category routing already
    # applied — sales_scoped is a real deterministic signal that
    # intent_explicit alone doesn't capture, and letting the LLM guess an
    # intent here previously clobbered it, e.g. a pricing question getting
    # silently re-routed to the Suggestions category). Never overrides
    # anything already resolved. Whatever it proposes is validated against
    # the real product/crop/intent lists before being trusted (see
    # llm_assisted_query_understanding). ──
    if not active_product and not active_crop and not intent_explicit and not sales_scoped and groq_api_key:
        llm_guess = llm_assisted_query_understanding(user_query, groq_api_key)
        if llm_guess:
            if llm_guess.get("product"):
                active_product = llm_guess["product"]
            if llm_guess.get("crop"):
                active_crop = llm_guess["crop"]
            if llm_guess.get("intent"):
                query_intent = llm_guess["intent"]
                intent_explicit = True
                category_filter = SUGGESTION_CATEGORY if query_intent == "suggestion" else category_filter

    retrieval_vector = query_vector
    retrieval_top_k = 100
    subject_for_embed = " ".join(filter(None, [active_crop, active_product]))
    if subject_for_embed:
        try:
            product_embed_response = pc.inference.embed(
                model="llama-text-embed-v2",
                inputs=[f"{subject_for_embed} product feedback sentiment complaints praise"],
                parameters={"input_type": "query", "dimension": EMBEDDING_DIMENSION}
            )
            retrieval_vector = product_embed_response[0].values
            retrieval_top_k = 300
        except Exception:
            retrieval_vector = query_vector

    # ── Comparison auto-detection ──
    periods = build_comparison_periods(all_months, all_years, all_weeks, index)

    # ── Relative time-window detection ("last 30 days", "last quarter",
    # "last 3 months") — only when the query didn't already pin an explicit
    # month/year (those always take priority) and isn't already a
    # multi-value comparison. Resolved BEFORE the ranking/trend branches so
    # those paths are time-scoped too; they used to return first, which is
    # why "the trend for the past three years" charted all of history. ──
    relative_window = None
    if not periods and not all_months and not all_years:
        window_desc = detect_relative_window(query_lower)
        if window_desc:
            relative_window = resolve_relative_window(window_desc, index)
    window_months = relative_window[0] if relative_window else None
    window_label = relative_window[1] if relative_window else None

    # ── Deterministic ranking path ──
    if aggregation_dimension:
        return _resolve_ranking(
            aggregation_dimension, query_intent, category_filter, detected_month, detected_year,
            index, output_format=output_format, groq_api_key=groq_api_key,
            window_months=window_months, window_label=window_label,
            top_n=requested_top_n or 10,
        )

    # ── Monthly trend path ──
    if wants_trend:
        return _resolve_trend(
            query_intent, category_filter, active_crop, active_product, index,
            output_format=output_format, groq_api_key=groq_api_key,
            window_months=window_months, window_label=window_label,
            detected_month=detected_month, detected_year=detected_year,
        )

    if periods:
        period_results = []
        for label, m, y, w in periods:
            p_pos, p_neg, p_neut = query_pinecone_for_timeframe(
                index, retrieval_vector, m, y, w, query_intent, top_k=retrieval_top_k, category_filter=category_filter
            )
            if active_product:
                p_pos = filter_bullets_by_product(p_pos, active_product)
                p_neg = filter_bullets_by_product(p_neg, active_product)
                p_neut = filter_bullets_by_product(p_neut, active_product)
            if active_crop:
                p_pos = filter_bullets_by_crop(p_pos, active_crop)
                p_neg = filter_bullets_by_crop(p_neg, active_crop)
                p_neut = filter_bullets_by_crop(p_neut, active_crop)

            MAX_BULLETS = 12
            period_results.append((label, p_pos[:MAX_BULLETS], p_neg[:MAX_BULLETS], p_neut[:MAX_BULLETS]))

        total_found = sum(len(pp) + len(pn) + len(pu) for _, pp, pn, pu in period_results)
        timeframe_label = " vs ".join(p[0] for p in periods)
        positive_bullets = negative_bullets = neutral_bullets = None
        target_year = None

    elif relative_window:
        window_months, window_label = relative_window
        # Collect per-month first, then interleave newest-month-first (see
        # _interleave_by_recency). Plain chronological concatenation here was
        # a real bug: downstream truncation takes the FIRST MAX_BULLETS
        # bullets, so "the last 6 months" was answered using only the OLDEST
        # month in the window and none of the recent ones.
        per_month_pos, per_month_neg, per_month_neut = [], [], []
        for w_month, w_year in window_months:
            w_pos, w_neg, w_neut = query_pinecone_for_timeframe(
                index, retrieval_vector, w_month, w_year, detected_week, query_intent, top_k=retrieval_top_k, category_filter=category_filter
            )
            per_month_pos.append(w_pos)
            per_month_neg.append(w_neg)
            per_month_neut.append(w_neut)

        positive_bullets = _interleave_by_recency(per_month_pos)
        negative_bullets = _interleave_by_recency(per_month_neg)
        neutral_bullets = _interleave_by_recency(per_month_neut)

        if active_product:
            positive_bullets = filter_bullets_by_product(positive_bullets, active_product)
            negative_bullets = filter_bullets_by_product(negative_bullets, active_product)
            neutral_bullets = filter_bullets_by_product(neutral_bullets, active_product)
        if active_crop:
            positive_bullets = filter_bullets_by_crop(positive_bullets, active_crop)
            negative_bullets = filter_bullets_by_crop(negative_bullets, active_crop)
            neutral_bullets = filter_bullets_by_crop(neutral_bullets, active_crop)

        total_found = len(positive_bullets) + len(negative_bullets) + len(neutral_bullets)
        timeframe_label = window_label
        target_year = None
        period_results = None

    else:
        target_year = detected_year
        fallback_triggered = False
        latest_index_year = None

        if detected_month and not target_year:
            latest_index_year = get_latest_year_from_index(index)
            target_year = latest_index_year

            pos, neg, neut = query_pinecone_for_timeframe(
                index, retrieval_vector, detected_month, target_year, detected_week, query_intent, top_k=retrieval_top_k, category_filter=category_filter
            )

            if (len(pos) + len(neg) + len(neut)) == 0:
                try:
                    fallback_year = str(int(latest_index_year) - 1)
                    pos_fb, neg_fb, neut_fb = query_pinecone_for_timeframe(
                        index, retrieval_vector, detected_month, fallback_year, detected_week, query_intent, top_k=retrieval_top_k, category_filter=category_filter
                    )
                    if (len(pos_fb) + len(neg_fb) + len(neut_fb)) > 0:
                        target_year = fallback_year
                        fallback_triggered = True
                except ValueError:
                    pass

        if wants_last_week and not detected_week:
            resolved_week = get_max_week_label(index, detected_month, target_year)
            if resolved_week:
                detected_week = resolved_week

        positive_bullets, negative_bullets, neutral_bullets = query_pinecone_for_timeframe(
            index, retrieval_vector, detected_month, target_year, detected_week, query_intent, top_k=retrieval_top_k, category_filter=category_filter
        )

        if active_product:
            positive_bullets = filter_bullets_by_product(positive_bullets, active_product)
            negative_bullets = filter_bullets_by_product(negative_bullets, active_product)
            neutral_bullets = filter_bullets_by_product(neutral_bullets, active_product)
        if active_crop:
            positive_bullets = filter_bullets_by_crop(positive_bullets, active_crop)
            negative_bullets = filter_bullets_by_crop(negative_bullets, active_crop)
            neutral_bullets = filter_bullets_by_crop(neutral_bullets, active_crop)

        total_found = len(positive_bullets) + len(negative_bullets) + len(neutral_bullets)

        timeframe_parts = []
        if detected_week:
            week_part = detected_week if "week" in detected_week.lower() else f"{detected_week} Week"
            timeframe_parts.append(week_part)
        if detected_month:
            timeframe_parts.append(f"of {detected_month}" if detected_week else detected_month)
        if target_year:
            timeframe_parts.append(target_year)
        timeframe_label = " ".join(timeframe_parts) or DEFAULT_TIMEFRAME_LABEL
        period_results = None

    header = build_header(query_intent, timeframe_label, active_product, periods, active_crop, category_filter=category_filter)
    badge = build_intent_badge(query_intent, active_product, periods, active_crop, category_filter=category_filter)
    subject_label = build_subject_label(active_product, active_crop)
    resolved_context = {"product": active_product, "crop": active_crop, "intent": query_intent}

    if total_found == 0:
        if periods:
            subject = f" for {subject_label}" if subject_label else ""
            reply = f"{badge}\n\n{header}No data found{subject} for the compared periods: {timeframe_label}."
        elif subject_label:
            suffix = f" in {timeframe_label}" if timeframe_label != DEFAULT_TIMEFRAME_LABEL else " in the ingested dataset"
            reply = f"{badge}\n\n{header}No data found for '{subject_label}'{suffix}."
        elif detected_month or detected_year or detected_week:
            reply = f"{badge}\n\n{header}No data found for {timeframe_label} in the ingested dataset."
        else:
            reply = (
                "I cannot generate this response. "
                "I am strictly locked to analyzed dataset metrics "
                "and cannot find relevant information for this query."
            )
        return {"kind": "no_data", "reply": reply, "context": resolved_context}

    # Breadth-seeking intents need a real sample to find genuine themes; a
    # narrow single-product lookup does not. One flat cap for both meant
    # "what is everyone talking about" was answered from the same 12
    # records as "what do growers think about Isabion".
    MAX_BULLETS = 40 if (query_intent == "topics" or (query_intent == "sentiment" and not subject_label)) else 12

    # True counts BEFORE truncation. These drive the chart, the KPIs and the
    # exported files — previously those were computed from the truncated
    # lists, so a month with 60 negative records exported "Negative: 12".
    true_counts = {
        "positive": len(positive_bullets) if positive_bullets is not None else 0,
        "negative": len(negative_bullets) if negative_bullets is not None else 0,
        "neutral": len(neutral_bullets) if neutral_bullets is not None else 0,
    }
    if periods and period_results:
        true_counts = {
            "positive": sum(len(pp) for _, pp, _, _ in period_results),
            "negative": sum(len(pn) for _, _, pn, _ in period_results),
            "neutral": sum(len(pu) for _, _, _, pu in period_results),
        }
    true_total = sum(true_counts.values())

    if periods:
        context_parts = []
        actual_point_count = 0
        for label, pos, neg, neut in period_results:
            section_lines = [f"=== {label} ==="]
            if query_intent == "complaint":
                pos = []
            elif query_intent == "positive":
                neg = []
                neut = []
            elif query_intent == "suggestion":
                pos = []
                neg = []
            if pos:
                section_lines.append("POSITIVE DATA:\n" + "\n".join(pos))
            if neg:
                section_lines.append("NEGATIVE DATA:\n" + "\n".join(neg))
            if neut:
                section_lines.append("OTHER DATA:\n" + "\n".join(neut))
            actual_point_count += len(pos) + len(neg) + len(neut)
            context_parts.append("\n".join(section_lines))
        combined_context = "\n\n".join(context_parts)
    else:
        if query_intent == "complaint":
            positive_bullets = []
            negative_bullets = negative_bullets[:MAX_BULLETS]
            neutral_bullets = neutral_bullets[:MAX_BULLETS]
        elif query_intent == "positive":
            positive_bullets = positive_bullets[:MAX_BULLETS]
            negative_bullets = []
            neutral_bullets = []
        elif query_intent == "suggestion":
            positive_bullets = []
            negative_bullets = []
            neutral_bullets = neutral_bullets[:MAX_BULLETS]
        else:
            positive_bullets = positive_bullets[:MAX_BULLETS]
            negative_bullets = negative_bullets[:MAX_BULLETS]
            neutral_bullets = neutral_bullets[:MAX_BULLETS]

        actual_point_count = len(positive_bullets) + len(negative_bullets) + len(neutral_bullets)

        context_parts = []
        if positive_bullets:
            context_parts.append("POSITIVE DATA:\n" + "\n".join(positive_bullets))
        if negative_bullets:
            context_parts.append("NEGATIVE DATA:\n" + "\n".join(negative_bullets))
        if neutral_bullets:
            context_parts.append("OTHER DATA:\n" + "\n".join(neutral_bullets))
        combined_context = "\n\n".join(context_parts)

    system_prompt = build_system_prompt(
        query_intent, timeframe_label, explicit_list_format, active_product, periods,
        active_crop=active_crop, output_format=output_format, wants_products_only=wants_products_only,
        avoid_repeat_text=avoid_repeat_text
    )
    # Be explicit with the model about sampling. Saying "N total" when N was
    # a truncated slice invited it to describe a 12-record sample as the
    # whole picture.
    if true_total > actual_point_count:
        volume_note = (
            f"Data Context — showing {actual_point_count} of {true_total} matching data points "
            f"(a sample, NOT the complete set). Base every statement only on the points below, "
            f"and do not state or imply totals, counts or proportions for the full dataset"
        )
    else:
        volume_note = (
            f"Data Context ({actual_point_count} distinct data point"
            f"{'s' if actual_point_count != 1 else ''} total — do not exceed this number)"
        )
    user_prompt = (
        f"Timeframe: {timeframe_label}\n\n"
        f"{volume_note}:\n{combined_context}\n\n"
        f"User Query: {user_query}"
    )
    response_token_budget = 900 if output_format in ("exec_summary", "table", "ppt") or query_intent == "topics" else 500

    return {
        "kind": "normal",
        "system_prompt": system_prompt,
        "user_prompt": user_prompt,
        "header": header,
        "badge": badge,
        "timeframe_label": timeframe_label,
        "subject_label": subject_label,
        "query_intent": query_intent,
        "output_format": output_format,
        "response_token_budget": response_token_budget,
        "periods": periods,
        "period_results": period_results,
        "positive_bullets": positive_bullets,
        "negative_bullets": negative_bullets,
        "neutral_bullets": neutral_bullets,
        "actual_point_count": actual_point_count,
        "true_counts": true_counts,
        "true_total": true_total,
        "context": resolved_context,
    }


def call_groq(system_prompt: str, user_prompt: str, groq_api_key: str, max_tokens: int = 500):
    """ Thin convenience wrapper — returns the Groq stream iterator. Each UI decides how to consume it: Streamlit updates st.empty() per chunk, a FastAPI endpoint would forward chunks as SSE. GROQ_MODEL is currently a reasoning model (openai/gpt-oss-20b) — reasoning_effort="low" keeps it from spending the whole token budget on internal chain-of-thought before ever emitting visible content. """
    client = Groq(api_key=groq_api_key)
    return client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.1,
        max_tokens=max_tokens,
        reasoning_effort="low",
        stream=True
    )


def _strip_code_fence(raw: str) -> str:
    """LLMs asked for raw JSON sometimes wrap it in a markdown code fence anyway — strip it before parsing."""
    return re.sub(r'^```(?:json)?\s*|\s*```$', '', raw.strip(), flags=re.MULTILINE).strip()


def generate_followup_suggestions(query_intent, subject_label, timeframe_label, full_response, groq_api_key, max_suggestions=3):
    """ One extra short Groq call after a normal response finishes: propose a few genuinely different follow-up questions the user could ask next, given what was just answered. Zero grounding risk — these are proposed QUESTIONS for the user to ask, not asserted facts, so there's nothing here for the LLM to hallucinate that could mislead anyone. Always returns a list (possibly empty) and never raises — a broken suggestion call must never break or block the main answer. """
    if not groq_api_key:
        return []
    try:
        client = Groq(api_key=groq_api_key)
        subject_bit = f" about {subject_label}" if subject_label else ""
        prompt = (
            f"A user just asked about {query_intent} feedback{subject_bit} for {timeframe_label}, "
            f"and received this answer:\n\n{full_response}\n\n"
            f"Suggest exactly {max_suggestions} short, genuinely different follow-up questions "
            f"they could ask next about this grower-feedback dataset (e.g. a different crop, "
            f"product, timeframe, or angle — not a rephrasing of the same question). "
            f"Return ONLY a JSON array of {max_suggestions} short question strings, nothing else — "
            f"no markdown, no explanation, no numbering."
        )
        resp = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=300,
            reasoning_effort="low",
        )
        parsed = json.loads(_strip_code_fence(resp.choices[0].message.content or ""))
        if not isinstance(parsed, list):
            return []
        return [str(s).strip() for s in parsed if str(s).strip()][:max_suggestions]
    except Exception:
        return []


def generate_deterministic_narrative(dimension_label, top_name, top_count, bullets, groq_api_key):
    """ Short LLM pass that adds 2-3 sentences of grounded color on top of an ALREADY-FINALIZED, deterministic ranking or trend result — called after the numbers are locked in, so it can only add narrative, never change a rank, a count, or a trend value. Grounded strictly in the real feedback bullets passed in; told explicitly to say "not clear" rather than guess if they don't explain the result. Returns "" on any failure or when there's nothing to ground it in — the numbers-only reply is always valid on its own without this. """
    if not groq_api_key or not bullets:
        return ""
    try:
        client = Groq(api_key=groq_api_key)
        context = "\n".join(f"- {b}" for b in bullets[:8])
        system_prompt = (
            "You are a data analyst. You must use ONLY the feedback bullets given below — "
            "never invent, assume, or add any detail not explicitly present in them. "
            "Write 2-3 sentences explaining the likely reason behind this result, grounded "
            "strictly in the bullets. If the bullets don't clearly explain why, say plainly "
            "that the data doesn't make the reason clear, instead of guessing."
        )
        user_prompt = f"Top result — {dimension_label}: {top_name} ({top_count} mentions)\n\nFeedback bullets:\n{context}"
        resp = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            temperature=0.2,
            max_tokens=250,
            reasoning_effort="low",
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception:
        return ""


def llm_assisted_query_understanding(user_query, groq_api_key):
    """ Fallback used ONLY when regex-based detection found NOTHING at all (no product, no crop, no explicit intent) — asks the LLM to propose a product/crop/intent for a genuinely ambiguous query. CRITICAL GUARDRAIL: whatever it proposes is validated against the real PRODUCT_LIST / CROP_LIST / the fixed intent enum before being trusted — anything not in those lists is silently discarded. This is exactly the check that would have caught an LLM inventing a product name that doesn't exist (observed live during testing) before it ever reached a user; the risk here is much lower than open-ended generation because the LLM can only ever "hit" a value that's already on an approved list, never introduce a new one. Returns None on any failure or when nothing usable survives validation. """
    if not groq_api_key:
        return None
    try:
        client = Groq(api_key=groq_api_key)
        system_prompt = (
            "You are a query classifier for a grower-feedback chatbot. Given a vague user "
            "question, propose the most likely product, crop, and intent it's about. "
            'Return ONLY a JSON object: {"product": "<name or null>", "crop": "<name or null>", '
            '"intent": "<one of: complaint, positive, suggestion, sentiment, topics>"}. '
            'Use "topics" when the question asks what people are talking about or '
            "discussing, rather than whether feedback is positive or negative. "
            "If you are not confident about a field, use null for it. Never invent a product "
            "or crop name — only propose one if it is a real product/crop explicitly present "
            "in the user's own question."
        )
        resp = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_query}],
            temperature=0.0,
            max_tokens=200,
            reasoning_effort="low",
        )
        parsed = json.loads(_strip_code_fence(resp.choices[0].message.content or ""))
        if not isinstance(parsed, dict):
            return None

        result = {"product": None, "crop": None, "intent": None}

        proposed_product = parsed.get("product")
        if isinstance(proposed_product, str) and proposed_product.strip().lower() in PRODUCT_LIST:
            result["product"] = proposed_product.strip().lower()

        proposed_crop = parsed.get("crop")
        if isinstance(proposed_crop, str) and proposed_crop.strip().lower() in CROP_LIST:
            result["crop"] = proposed_crop.strip().lower()

        proposed_intent = parsed.get("intent")
        if proposed_intent in ("complaint", "positive", "suggestion", "sentiment", "topics"):
            result["intent"] = proposed_intent

        if not any(result.values()):
            return None
        return result
    except Exception:
        return None


def finalize_normal_response(state: dict, full_response: str) -> dict:
    """ Second phase for kind="normal": once the caller has streamed the full LLM response text, compute the chart, export rows, and ready-to-serve CSV/Excel/PPTX bytes. Returns {"final_reply", "chart", "downloads"}. """
    periods = state["periods"]
    period_results = state["period_results"]
    positive_bullets = state["positive_bullets"]
    negative_bullets = state["negative_bullets"]
    neutral_bullets = state["neutral_bullets"]
    timeframe_label = state["timeframe_label"]
    header = state["header"]
    badge = state["badge"]
    subject_label = state["subject_label"]
    query_intent = state["query_intent"]
    actual_point_count = state["actual_point_count"]

    export_rows = []
    if periods:
        for label, pp, pn, pu in period_results:
            for b in pp:
                export_rows.append({"Period": label, "Sentiment": "Positive", "Feedback": b})
            for b in pn:
                export_rows.append({"Period": label, "Sentiment": "Negative", "Feedback": b})
            for b in pu:
                export_rows.append({"Period": label, "Sentiment": "Other", "Feedback": b})
    else:
        for b in positive_bullets:
            export_rows.append({"Period": timeframe_label, "Sentiment": "Positive", "Feedback": b})
        for b in negative_bullets:
            export_rows.append({"Period": timeframe_label, "Sentiment": "Negative", "Feedback": b})
        for b in neutral_bullets:
            export_rows.append({"Period": timeframe_label, "Sentiment": "Other", "Feedback": b})

    # Charts and KPIs report the TRUE match counts, not the truncated slice
    # that was sent to the model. Deriving them from the truncated lists
    # meant a month with 60 negative records rendered as "Negative: 12" —
    # in the PowerPoint and Excel files a business user downloads.
    true_counts = state.get("true_counts") or {
        "positive": len(positive_bullets or []),
        "negative": len(negative_bullets or []),
        "neutral": len(neutral_bullets or []),
    }
    true_total = state.get("true_total", sum(true_counts.values()))

    if periods:
        chart_labels = [label for label, *_ in period_results]
        chart_values = [len(pp) + len(pn) + len(pu) for _, pp, pn, pu in period_results]
        chart_title = "Total Data Points by Period"
    else:
        chart_labels = ["Positive", "Negative", "Other"]
        chart_values = [true_counts["positive"], true_counts["negative"], true_counts["neutral"]]
        chart_title = "Sentiment Breakdown"

    kpis = {
        "Total matching records": true_total,
        "Positive": true_counts["positive"],
        "Negative": true_counts["negative"],
        "Other": true_counts["neutral"],
    }
    if true_total > actual_point_count:
        kpis["Shown to the model"] = f"{actual_point_count} (sample)"
    summary_points = split_into_points(full_response, max_points=6)

    pptx_bytes = build_pptx_report(
        title=subject_label or f"{query_intent.title()} Analysis",
        subtitle=timeframe_label,
        exec_summary_lines=summary_points[:4] or ["No summary available."],
        kpis=kpis,
        chart_title=chart_title,
        chart_labels=chart_labels,
        chart_values=chart_values,
        chart_type="column",
        table_headers=["Sentiment", "Feedback"],
        table_rows=[(row["Sentiment"], row["Feedback"]) for row in export_rows],
        insights=summary_points[:4],
        recommendations=summary_points[4:6],
    )

    downloads = {
        "csv": build_csv_bytes([{"Label": l, "Value": v} for l, v in zip(chart_labels, chart_values)]),
        "excel": build_excel_bytes(export_rows) if export_rows else None,
        "pptx": pptx_bytes,
    }

    final_reply = badge + "\n\n" + header + full_response

    # Chart is only shown inline in the chat when the user explicitly asked
    # for one — downloads (CSV/Excel/PPTX) keep their chart regardless,
    # since those are opt-in via an explicit click, not shown unprompted.
    wants_chart = state.get("output_format") == "chart"

    return {
        "final_reply": final_reply,
        "chart": {"type": "bar", "title": chart_title, "labels": chart_labels, "values": chart_values} if wants_chart else None,
        "downloads": downloads,
        "kpis": kpis,
    }


def _resolve_ranking(aggregation_dimension, query_intent, category_filter, detected_month, detected_year, index, output_format=None, groq_api_key=None, window_months=None, window_label=None, top_n=10) -> dict:
    agg_filter = {}
    # Time scoping now reaches this path. Previously ranking returned before
    # relative-window resolution ran and never received a month/year, so
    # "which crop had the most complaints last quarter" silently ranked over
    # all history and labelled itself "all available data".
    if window_months:
        agg_filter["month"] = {"$in": sorted({m for m, _ in window_months})}
        agg_filter["year"] = {"$in": sorted({y for _, y in window_months})}
    else:
        if detected_month:
            agg_filter["month"] = {"$eq": detected_month}
        if detected_year:
            agg_filter["year"] = {"$eq": detected_year}
    if query_intent == "positive":
        agg_filter["sentiment"] = {"$eq": "positive"}
    elif query_intent == "complaint":
        agg_filter["sentiment"] = {"$eq": "negative"}
    elif category_filter:
        agg_filter["category"] = {"$eq": category_filter}

    agg_matches, complete = fetch_matches_for_aggregation(index, agg_filter)
    field = "crop" if aggregation_dimension == "crop" else "products"
    ranking = rank_by_field(agg_matches, field, top_n=top_n)

    badge = f"📊 {aggregation_dimension.title()} Ranking"
    if window_label:
        scope_label = window_label
    else:
        scope_bits = [b for b in (detected_month, detected_year) if b]
        scope_label = " ".join(scope_bits) if scope_bits else "all available data"
    header = f"📊 {aggregation_dimension.title()}-wise Ranking ({scope_label}):\n\n"

    if not ranking:
        reply = (
            f"{header}No {aggregation_dimension} tags were found in the matched "
            f"records for {scope_label} — nothing to rank."
        )
        return {"kind": "ranking", "reply": reply, "badge": badge, "chart": None, "downloads": None}

    table_lines = ["| Rank | " + aggregation_dimension.title() + " | Mentions |", "|---|---|---|"]
    for i, (name, count) in enumerate(ranking, start=1):
        table_lines.append(f"| {i} | {name} | {count} |")
    top_name, top_count = ranking[0]
    tagged_total = sum(c for _, c in rank_by_field(agg_matches, field, top_n=10_000))
    # Say "at least" when the fetch was capped, rather than presenting a
    # truncated subset as a complete census.
    basis = (
        f"Based on {len(agg_matches)} matched records ({tagged_total} tagged mentions)"
        if complete else
        f"Based on a sample of {len(agg_matches)} matched records — more exist, so these counts are lower bounds"
    )
    reply = (
        f"{header}"
        f"{basis}, **{top_name}** ranks highest "
        f"with {top_count} mention{'s' if top_count != 1 else ''}.\n\n"
        + "\n".join(table_lines)
    )

    # Optional narrative color on top of the already-finalized numbers above
    # — the ranking/table/counts are already fixed by this point, so this
    # can only add explanation, never change what was already computed.
    top_bullets = []
    top_name_lower = top_name.lower()
    for m in agg_matches:
        md = m.get("metadata", {})
        if top_name_lower in str(md.get(field, "")).lower():
            v = str(md.get("value", "")).strip()
            if v and v not in top_bullets:
                top_bullets.append(v)
        if len(top_bullets) >= 8:
            break
    narrative = generate_deterministic_narrative(aggregation_dimension, top_name, top_count, top_bullets, groq_api_key)
    if narrative:
        reply += f"\n\n*Why {top_name} likely ranks highest:* {narrative}"

    labels = [n for n, _ in ranking]
    values = [c for _, c in ranking]

    pptx_bytes = build_pptx_report(
        title=f"{aggregation_dimension.title()}-wise Ranking",
        subtitle=scope_label,
        exec_summary_lines=[f"Based on {len(agg_matches)} matched records, {top_name} ranks highest with {top_count} mentions."],
        kpis={"Total matched records": len(agg_matches), "Top " + aggregation_dimension: top_name, "Top mentions": top_count},
        chart_title=f"{aggregation_dimension.title()} Mentions", chart_labels=labels, chart_values=values,
        chart_type="column",
        table_headers=["Rank", aggregation_dimension.title(), "Mentions"],
        table_rows=[(i, n, c) for i, (n, c) in enumerate(ranking, start=1)],
        insights=[f"{n} — {c} mentions" for n, c in ranking[:5]],
    )
    downloads = {
        "csv": build_csv_bytes([{aggregation_dimension.title(): n, "Mentions": c} for n, c in ranking]),
        "excel": build_excel_bytes([{aggregation_dimension.title(): n, "Mentions": c} for n, c in ranking]),
        "pptx": pptx_bytes,
    }

    wants_chart = output_format == "chart"
    return {
        "kind": "ranking",
        "reply": reply,
        "badge": badge,
        "chart": {"type": "bar", "title": f"{aggregation_dimension.title()} Mentions", "labels": labels, "values": values} if wants_chart else None,
        "downloads": downloads,
    }


def _resolve_trend(query_intent, category_filter, active_crop, active_product, index, output_format=None, groq_api_key=None, window_months=None, window_label=None, detected_month=None, detected_year=None) -> dict:
    trend_filter = {}
    # Trend previously received no time parameters at all, so "the monthly
    # complaint trend for the past three years" charted all of history and
    # showed no timeframe in its header.
    if window_months:
        trend_filter["month"] = {"$in": sorted({m for m, _ in window_months})}
        trend_filter["year"] = {"$in": sorted({y for _, y in window_months})}
    elif detected_year:
        trend_filter["year"] = {"$eq": detected_year}
        if detected_month:
            trend_filter["month"] = {"$eq": detected_month}

    if query_intent == "positive":
        trend_filter["sentiment"] = {"$eq": "positive"}
    elif query_intent == "complaint":
        trend_filter["sentiment"] = {"$eq": "negative"}
    elif category_filter:
        trend_filter["category"] = {"$eq": category_filter}

    trend_matches, complete = fetch_matches_for_aggregation(index, trend_filter)
    if active_crop:
        trend_matches = [m for m in trend_matches if _mentions(str(m.get("metadata", {}).get("value", "")), active_crop)]
    if active_product:
        trend_matches = [m for m in trend_matches if _mentions(str(m.get("metadata", {}).get("value", "")), active_product)]

    monthly_counts = densify_monthly_counts(compute_monthly_trend(trend_matches))
    growth_series = compute_growth_series(monthly_counts)

    subject_label = build_subject_label(active_product, active_crop)
    subject_bit = f" for {subject_label}" if subject_label else ""
    scope_bit = f" ({window_label})" if window_label else (f" ({detected_year})" if detected_year else "")
    header = f"📈 Monthly Trend Analysis{subject_bit}{scope_bit}:\n\n"

    if not monthly_counts:
        reply = f"{header}No dated records were found{subject_bit} to build a monthly trend from."
        return {"kind": "no_data", "reply": reply}

    highest = max(monthly_counts, key=lambda x: x[1])
    lowest = min(monthly_counts, key=lambda x: x[1])
    table_lines = ["| Month | Count | MoM Growth |", "|---|---|---|"]
    for (label, count), growth in zip(monthly_counts, growth_series):
        growth_str = "—" if growth is None else f"{'+' if growth >= 0 else ''}{growth}%"
        table_lines.append(f"| {label} | {count} | {growth_str} |")

    reply = (
        f"{header}"
        f"Highest month: **{highest[0]}** ({highest[1]} records). "
        f"Lowest month: **{lowest[0]}** ({lowest[1]} records).\n\n"
        + "\n".join(table_lines)
    )

    # Optional narrative color on top of the already-finalized numbers above
    # — grounded in real bullets from the highest month only, and generated
    # after that month is already fixed, so it can only add explanation.
    highest_month, highest_year = highest[0].split(" ", 1)
    highest_bullets = []
    for m in trend_matches:
        md = m.get("metadata", {})
        if md.get("month") == highest_month and md.get("year") == highest_year:
            v = str(md.get("value", "")).strip()
            if v and v not in highest_bullets:
                highest_bullets.append(v)
        if len(highest_bullets) >= 8:
            break
    narrative = generate_deterministic_narrative("month", highest[0], highest[1], highest_bullets, groq_api_key)
    if narrative:
        reply += f"\n\n*Why {highest[0]} was the strongest month:* {narrative}"

    labels = [l for l, _ in monthly_counts]
    values = [c for _, c in monthly_counts]

    pptx_bytes = build_pptx_report(
        title=f"Monthly Trend Analysis{subject_bit}",
        subtitle=f"{monthly_counts[0][0]} – {monthly_counts[-1][0]}",
        exec_summary_lines=[
            f"Highest month: {highest[0]} with {highest[1]} records.",
            f"Lowest month: {lowest[0]} with {lowest[1]} records.",
        ],
        kpis={"Total records": sum(values), "Highest month": f"{highest[0]} ({highest[1]})",
              "Lowest month": f"{lowest[0]} ({lowest[1]})"},
        chart_title="Monthly Trend", chart_labels=labels, chart_values=values, chart_type="line",
        table_headers=["Month", "Count", "MoM Growth %"],
        table_rows=[(l, c, ("—" if g is None else f"{g}%")) for (l, c), g in zip(monthly_counts, growth_series)],
        insights=[f"{highest[0]} was the strongest month ({highest[1]} records).",
                  f"{lowest[0]} was the weakest month ({lowest[1]} records)."],
    )
    downloads = {
        "csv": build_csv_bytes([{"Month": l, "Count": c, "MoM Growth %": g} for (l, c), g in zip(monthly_counts, growth_series)]),
        "excel": build_excel_bytes([{"Month": l, "Count": c, "MoM Growth %": g} for (l, c), g in zip(monthly_counts, growth_series)]),
        "pptx": pptx_bytes,
    }

    wants_chart = output_format == "chart"
    return {
        "kind": "trend",
        "reply": reply,
        "badge": "📈 Monthly Trend",
        "chart": {"type": "line", "title": "Monthly Trend", "labels": labels, "values": values} if wants_chart else None,
        "downloads": downloads,
    }
