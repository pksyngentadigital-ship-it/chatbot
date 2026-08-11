"""
Header/badge/prompt text builders, and the PPTX/Excel/CSV export builders.
"""
import pandas as pd
from io import BytesIO
from pptx import Presentation

import vog_core as vc


# ── build_subject_label ──

def test_subject_label_combines_crop_and_product():
    assert vc.build_subject_label("isabion", "wheat") == "Wheat + Isabion"


def test_subject_label_collapses_when_crop_equals_product():
    assert vc.build_subject_label("wheat", "wheat") == "Wheat"


def test_subject_label_none_when_neither_present():
    assert vc.build_subject_label(None, None) is None


# ── build_header: product/crop must always take priority over generic period heading ──

def test_header_uses_subject_over_generic_period():
    header = vc.build_header("complaint", "January 2026", "isabion", [], active_crop=None)
    assert "Isabion" in header
    assert "Complaints about" in header


def test_header_generic_when_no_subject():
    header = vc.build_header("sentiment", "January 2026", None, [])
    assert "Sentiments of January 2026" in header


def test_header_comparison_is_labeled_as_comparison():
    periods = [("January 2026", "January", "2026", None), ("February 2026", "February", "2026", None)]
    header = vc.build_header("complaint", "n/a", None, periods)
    assert "Comparison" in header
    assert "January 2026" in header and "February 2026" in header


# ── build_intent_badge ──

def test_badge_shows_subject_when_present():
    badge = vc.build_intent_badge("sentiment", "isabion", [], active_crop=None)
    assert "Isabion" in badge


def test_badge_comparison_label():
    periods = [("A", None, None, None), ("B", None, None, None)]
    assert vc.build_intent_badge("sentiment", None, periods) == "🔀 Comparison"


def test_badge_default_sentiment_overview():
    assert vc.build_intent_badge("sentiment", None, []) == "🌾 Sentiment Overview"


# ── build_system_prompt: product-only guardrail must exclude diseases/pests ──

def test_system_prompt_products_only_excludes_diseases():
    prompt = vc.build_system_prompt(
        "sentiment", "January 2026", explicit_list_format=False,
        active_product=None, periods=[], wants_products_only=True,
    )
    assert "Do NOT list crop diseases" in prompt


def test_system_prompt_omits_products_only_clause_when_not_requested():
    prompt = vc.build_system_prompt(
        "sentiment", "January 2026", explicit_list_format=False,
        active_product=None, periods=[], wants_products_only=False,
    )
    assert "Do NOT list crop diseases" not in prompt


def test_system_prompt_focuses_exclusively_on_active_product():
    prompt = vc.build_system_prompt(
        "sentiment", "January 2026", explicit_list_format=False,
        active_product="isabion", periods=[],
    )
    assert "Focus EXCLUSIVELY on the product 'Isabion'" in prompt


def test_system_prompt_table_output_format_overrides_bullet_format():
    prompt = vc.build_system_prompt(
        "sentiment", "January 2026", explicit_list_format=True,
        active_product=None, periods=[], output_format="table",
    )
    assert "markdown table" in prompt
    assert "real markdown bullet list" not in prompt


# ── Export builders: must produce valid, non-empty files ──

def test_build_csv_bytes_round_trips_through_pandas():
    rows = [{"Label": "Positive", "Value": 3}, {"Label": "Negative", "Value": 1}]
    csv_bytes = vc.build_csv_bytes(rows)
    df = pd.read_csv(BytesIO(csv_bytes))
    assert list(df["Label"]) == ["Positive", "Negative"]
    assert list(df["Value"]) == [3, 1]


def test_build_excel_bytes_round_trips_through_pandas():
    rows = [{"Period": "Jan 2026", "Sentiment": "Positive", "Feedback": "Great results"}]
    xlsx_bytes = vc.build_excel_bytes(rows)
    df = pd.read_excel(BytesIO(xlsx_bytes))
    assert df.iloc[0]["Feedback"] == "Great results"


def test_build_pptx_report_produces_valid_deck_with_expected_slides():
    pptx_bytes = vc.build_pptx_report(
        title="Test Report", subtitle="January 2026",
        exec_summary_lines=["Everything looks good."],
        kpis={"Total": 10},
        chart_title="Sentiment Breakdown",
        chart_labels=["Positive", "Negative"],
        chart_values=[7, 3],
        table_headers=["Sentiment", "Feedback"],
        table_rows=[("Positive", "Great product")],
        insights=["Insight one"],
        recommendations=["Do this next"],
    )
    assert isinstance(pptx_bytes, bytes)
    assert len(pptx_bytes) > 0
    prs = Presentation(BytesIO(pptx_bytes))
    # Title, Executive Summary, KPIs, Chart, Table, Insights = 6 slides
    assert len(prs.slides) == 6


def test_build_pptx_report_skips_optional_slides_when_no_data():
    pptx_bytes = vc.build_pptx_report(
        title="Empty Report", subtitle="n/a",
        exec_summary_lines=[], kpis={}, chart_title="", chart_labels=[], chart_values=[],
    )
    prs = Presentation(BytesIO(pptx_bytes))
    # No chart/table/insights data -> only Title, Exec Summary, KPIs
    assert len(prs.slides) == 3
